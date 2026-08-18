"""Genere la PRESENTATION PowerPoint de la partie LLM (16:9).

Reprend la meme structure, les memes tableaux et la meme logique que le rapport
Word (docs/generate_rapport_llm.py) : 3 parties (chunking, architectures,
optimisations) + synthese, verdicts colores vert/rouge, analyses "pourquoi",
comparaisons cote a cote, sources academiques verifiees (liens cliquables).
Aucun em-dash ni fleche dans le texte.
"""

from __future__ import annotations
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# Palette (sobre, identique au rapport : bleu ardoise + accent cuivre)
BLEU = RGBColor(0x1F, 0x3A, 0x5F)
BLEU2 = RGBColor(0x2E, 0x4D, 0x6B)
GRIS = RGBColor(0x5A, 0x5A, 0x5A)
NOIR = RGBColor(0x22, 0x22, 0x22)
VERT = RGBColor(0x2E, 0x6B, 0x3F)
ROUGE = RGBColor(0x9E, 0x2B, 0x25)
ACCENT = RGBColor(0xB5, 0x65, 0x1D)
BLANC = RGBColor(0xFF, 0xFF, 0xFF)
HDR_BG = RGBColor(0x1F, 0x3A, 0x5F)
ZEBRA = RGBColor(0xF4, 0xF6, 0xF9)
VERT_BG = RGBColor(0xEA, 0xF1, 0xEC)
ROUGE_BG = RGBColor(0xF7, 0xEC, 0xEB)
BLEU_BG = RGBColor(0xEE, 0xF2, 0xF6)
JAUNE_BG = RGBColor(0xFB, 0xF4, 0xEA)
GRIS_BG = RGBColor(0xF4, 0xF5, 0xF7)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W = prs.slide_width
H = prs.slide_height
CL = Inches(0.5)          # marge gauche contenu
CW = Inches(12.33)        # largeur contenu


# --------------------------------------------------------------------------
# Helpers
# --------------------------------------------------------------------------
def _no_line(shape):
    shape.line.fill.background()
    shape.shadow.inherit = False


def _table_nogrid(table):
    """Applique le style 'No Style, No Grid' : aucune ligne (surtout verticale)."""
    tblPr = table._tbl.find(qn("a:tblPr"))
    if tblPr is None:
        return
    for attr in ("firstRow", "lastRow", "firstCol", "lastCol", "bandRow", "bandCol"):
        tblPr.set(attr, "0")
    sid = tblPr.find(qn("a:tableStyleId"))
    if sid is None:
        sid = tblPr.makeelement(qn("a:tableStyleId"), {})
        tblPr.append(sid)
    sid.text = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"


def _runs(tf, parts, size=12, color=NOIR, align=PP_ALIGN.LEFT, space_after=2):
    """parts = liste de (texte, bold, color) pour un paragraphe."""
    p = tf.paragraphs[0] if not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    for (txt, bold, col) in parts:
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = col if col else color
    return p


def blank():
    return prs.slides.add_slide(BLANK)


def banner(slide, title):
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.95))
    band.fill.solid()
    band.fill.fore_color.rgb = BLEU
    _no_line(band)
    tf = band.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.4)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Georgia"
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = BLANC
    return slide


def content_slide(title):
    s = blank()
    banner(s, title)
    return s


def textbox(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def paragraph(slide, left, top, width, text, size=13, italic=False, bold=False, color=NOIR):
    tf = textbox(slide, left, top, width, Inches(0.6))
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.italic = italic
    r.font.bold = bold
    r.font.color.rgb = color
    return tf


def bullets(slide, left, top, width, height, items, size=14):
    """items = liste de (prefixe_gras, texte) ou (None, texte)."""
    tf = textbox(slide, left, top, width, height)
    for i, it in enumerate(items):
        pref, txt = it
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        rb = p.add_run()
        rb.text = "•  "
        rb.font.size = Pt(size)
        rb.font.color.rgb = BLEU2
        if pref:
            r = p.add_run()
            r.text = pref
            r.font.bold = True
            r.font.size = Pt(size)
            r.font.color.rgb = NOIR
        r2 = p.add_run()
        r2.text = txt
        r2.font.size = Pt(size)
        r2.font.color.rgb = NOIR
    return tf


def box(slide, left, top, width, height, label, text, fill, label_color, size=12):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = label_color
    sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    p = tf.paragraphs[0]
    if label:
        r = p.add_run()
        r.text = label + "  "
        r.font.bold = True
        r.font.size = Pt(size)
        r.font.color.rgb = label_color
    r2 = p.add_run()
    r2.text = text
    r2.font.size = Pt(size)
    r2.font.color.rgb = NOIR
    return sh


def verdict(slide, left, top, width, height, text, size=12):
    low = text.lower()
    if low.startswith(("adopt", "retenu", "prometteur")):
        box(slide, left, top, width, height, "Verdict", text, VERT_BG, VERT, size)
    elif low.startswith(("rejet", "ecart", "écart", "non concl", "abandonn")):
        box(slide, left, top, width, height, "Verdict", text, ROUGE_BG, ROUGE, size)
    else:
        box(slide, left, top, width, height, "Verdict", text, BLEU_BG, BLEU, size)


def analyse(slide, left, top, width, height, text, size=12):
    box(slide, left, top, width, height, "Analyse", text, BLEU_BG, BLEU, size)


def note(slide, left, top, width, height, text, size=12):
    box(slide, left, top, width, height, "Note", text, JAUNE_BG, RGBColor(0x8A, 0x6D, 0x00), size)


def _set_cell(cell, text, size=10, bold=False, color=NOIR, fill=None, align=PP_ALIGN.LEFT):
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.06)
    cell.margin_right = Inches(0.06)
    cell.margin_top = Inches(0.02)
    cell.margin_bottom = Inches(0.02)
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = str(text)
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    else:
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLANC


def table(slide, left, top, width, headers, rows, col_fracs=None,
          green=None, red=None, font=10, row_h=0.34):
    green = green or set()
    red = red or set()
    NEG = ("plus lent", "plus lente", "incomplet", "plante", "vide", "rejet")
    POS = ("adopt", "retenu")
    nrows = len(rows) + 1
    ncols = len(headers)
    height = Inches(row_h * nrows)
    gf = slide.shapes.add_table(nrows, ncols, left, top, width, height)
    t = gf.table
    t.first_row = False
    t.horz_banding = False
    _table_nogrid(t)   # supprime la grille (et donc toute ligne verticale)
    if col_fracs:
        total = sum(col_fracs)
        for j, fr in enumerate(col_fracs):
            t.columns[j].width = Emu(int(int(width) * fr / total))
    for j, h in enumerate(headers):
        _set_cell(t.cell(0, j), h, size=font, bold=True, color=BLANC, fill=HDR_BG)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            sval = str(val)
            low = sval.lower()
            color, bold = NOIR, False
            if (i, j) in green:
                color, bold = VERT, True
            elif (i, j) in red:
                color, bold = ROUGE, True
            elif sval.startswith("-") and "%" in sval:
                color, bold = VERT, True
            elif any(k in low for k in NEG):
                color, bold = ROUGE, True
            elif any(low.startswith(k) for k in POS):
                color, bold = VERT, True
            fill = ZEBRA if i % 2 == 1 else BLANC
            _set_cell(t.cell(i + 1, j), sval, size=font, bold=bold, color=color, fill=fill)
    for i in range(nrows):
        t.rows[i].height = Inches(row_h)
    return gf


def two_col(slide, top, titleA, textA, srcA, titleB, textB, srcB,
            sentA="neutre", sentB="neutre", height=Inches(3.0)):
    fills = {"bon": VERT_BG, "mauvais": ROUGE_BG, "neutre": BLEU_BG}
    cols = {"bon": VERT, "mauvais": ROUGE, "neutre": BLEU}
    w = Inches(5.95)
    gap = Inches(0.43)
    for (left, ti, tx, src, sent) in [
        (CL, titleA, textA, srcA, sentA),
        (CL + w + gap, titleB, textB, srcB, sentB),
    ]:
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, height)
        sh.fill.solid()
        sh.fill.fore_color.rgb = fills[sent]
        _no_line(sh)
        tf = sh.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)
        tf.margin_top = Inches(0.1)
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = ti
        r.font.bold = True
        r.font.size = Pt(13)
        r.font.color.rgb = cols[sent]
        p2 = tf.add_paragraph()
        p2.space_before = Pt(4)
        r2 = p2.add_run()
        r2.text = tx
        r2.font.size = Pt(10.5)
        r2.font.italic = True
        r2.font.color.rgb = NOIR
        p3 = tf.add_paragraph()
        p3.space_before = Pt(4)
        r3 = p3.add_run()
        r3.text = "Source : " + src
        r3.font.size = Pt(8.5)
        r3.font.color.rgb = BLEU2


def sources(slide, left, top, width, items, height=Inches(0.5)):
    tf = textbox(slide, left, top, width, height)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Source(s) : "
    r.font.bold = True
    r.font.size = Pt(9)
    r.font.color.rgb = GRIS
    for i, (label, url) in enumerate(items):
        if i > 0:
            s = p.add_run()
            s.text = "      "
            s.font.size = Pt(9)
        r = p.add_run()
        r.text = label
        r.font.size = Pt(9)
        r.font.color.rgb = BLEU
        r.font.underline = True
        r.hyperlink.address = url


def excerpt(slide, left, top, width, height, lines, source, size=10):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = GRIS_BG
    _no_line(sh)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.12)
    tf.margin_top = Inches(0.06)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(size)
        r.font.italic = True
        r.font.color.rgb = NOIR
    p = tf.add_paragraph()
    r = p.add_run()
    r.text = "Extrait de : " + source
    r.font.size = Pt(8.5)
    r.font.color.rgb = BLEU2


def divider(num, title, sub=None):
    s = blank()
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BLEU
    _no_line(bg)
    tf = bg.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(1.0)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = num
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0x9D, 0xC3, 0xE6)
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = title
    r2.font.name = "Georgia"
    r2.font.size = Pt(40)
    r2.font.bold = True
    r2.font.color.rgb = BLANC
    if sub:
        p3 = tf.add_paragraph()
        p3.space_before = Pt(10)
        r3 = p3.add_run()
        r3.text = sub
        r3.font.size = Pt(16)
        r3.font.color.rgb = RGBColor(0xD6, 0xE4, 0xF0)
    return s


# ==========================================================================
# SLIDE 1 — TITRE
# ==========================================================================
s = blank()
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
bg.fill.solid()
bg.fill.fore_color.rgb = BLEU
_no_line(bg)
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.6), W, Inches(2.4))
band.fill.solid()
band.fill.fore_color.rgb = RGBColor(0x16, 0x3A, 0x5C)
_no_line(band)
tf = band.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tf.margin_left = Inches(0.9)
tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run()
r.text = "Génération de comptes rendus de réunion par LLM local"
r.font.name = "Georgia"
r.font.size = Pt(34)
r.font.bold = True
r.font.color.rgb = BLANC
p2 = tf.add_paragraph()
p2.space_before = Pt(10)
r2 = p2.add_run()
r2.text = "Du découpage du transcript à l'architecture agentique : architectures, optimisations et benchmarks"
r2.font.size = Pt(16)
r2.font.color.rgb = RGBColor(0x9D, 0xC3, 0xE6)
foot = textbox(s, Inches(0.9), Inches(5.3), Inches(11), Inches(0.8))
p = foot.paragraphs[0]
r = p.add_run()
r.text = "Yele Consulting, rapport technique (partie LLM)   |   Expériences mars à juin 2026"
r.font.size = Pt(13)
r.font.color.rgb = RGBColor(0xD6, 0xE4, 0xF0)

# ==========================================================================
# SLIDE 2 — CONTEXTE
# ==========================================================================
s = content_slide("0. Contexte, objectif et méthode")
paragraph(s, CL, Inches(1.15), CW,
          "Objectif : transformer le transcript d'une réunion (diarisation et transcription "
          "locales) en compte rendu structuré, sous trois contraintes non négociables.",
          size=13)
bullets(s, CL, Inches(1.85), CW, Inches(1.6), [
    ("Contrainte 1 : ", "100 % local, CPU uniquement (8 à 16 Go de RAM), sans GPU ni cloud."),
    ("Contrainte 2 : ", "zéro hallucination, ne restituer que ce qui est explicitement dit."),
    ("Contrainte 3 : ", "fonctionner en live comme sur un audio importé."),
], size=14)
box(s, CL, Inches(3.5), CW, Inches(1.0), "Réunion de référence",
    "Toutes les expériences et tous les extraits portent sur LA MÊME réunion RTE (environ 1 h, "
    "environ 9 participants dont Jérôme Picault et Maya Sahraoui, IA générative et orchestration MCP). "
    "Transcripts : transcript1 (556 seg), transcript_formatted (247 seg), dicte_audio_3 (266 seg).",
    BLEU_BG, BLEU, size=12)
note(s, CL, Inches(4.7), CW, Inches(1.5),
     "Pas de gold standard humain dans le projet : les comparaisons portent sur des métriques "
     "objectives (durée, RAM, nombre d'appels, items extraits) et une évaluation qualitative "
     "manuelle. Tout chiffre vient verbatim des metrics.json et des logs. Les CR complets cités "
     "sont dans docs/comptes_rendus_references/.", size=12)

# ==========================================================================
# PARTIE 1
# ==========================================================================
divider("Partie 1", "Le découpage (chunking)",
        "Découper 1 h de réunion en morceaux cohérents traités un par un")

# 1.1 approches V1/V2
s = content_slide("1.1  Approches de découpage : clustering HDBSCAN (V1 puis V2)")
paragraph(s, CL, Inches(1.1), CW,
          "V1 : fenêtres glissantes, embeddings MiniLM, clustering HDBSCAN, 1 LLM par cluster. "
          "Défaut : clusters non chronologiques, 30 à 40 % du contenu jeté comme bruit.", size=13)
paragraph(s, CL, Inches(1.95), CW, "V2 : HDBSCAN amélioré (mesuré sur transcript1, 556 seg).",
          size=13, bold=True)
table(s, CL, Inches(2.45), Inches(8.5),
      ["Modèle", "Durée", "RAM serveur", "Clusters"],
      [["Ministral 3B", "39 min 22", "3 786 Mo", "5"],
       ["Qwen2.5 3B", "20 min 48", "2 670 Mo", "5"],
       ["Qwen3 4B", "48 min 55", "4 670 Mo", "5"],
       ["SmolLM3 3B", "24 min 48", "3 870 Mo", "5"]],
      col_fracs=[3, 2, 2, 1.5], font=11)
verdict(s, CL, Inches(5.4), CW, Inches(0.9),
        "Abandonné à terme. Plus robuste que V1 mais toujours dépendant du clustering, et trop lent.")

# 1.1 V3 boundary
s = content_slide("1.1  Détection de frontières sémantiques (V3, RETENUE)")
paragraph(s, CL, Inches(1.1), CW,
          "Frontières par chute de similarité cosine entre fenêtres (lissage sigma=2, vallées au "
          "percentile 5). Chunks chronologiques, 1 LLM par chunk (transcript1, 556 seg, 9 chunks).", size=13)
table(s, CL, Inches(2.0), Inches(9),
      ["Modèle", "V2 (HDBSCAN)", "V3 (boundary)", "Gain"],
      [["Ministral 3B", "39 min 22", "26 min 37", "-32 %"],
       ["Qwen2.5 3B", "20 min 48", "12 min 13", "-41 %"],
       ["Qwen3 4B", "48 min 55", "32 min 40", "-33 %"],
       ["SmolLM3 3B", "24 min 48", "14 min 56", "-40 %"]],
      col_fracs=[2.2, 2, 2, 1.2], font=11)
verdict(s, CL, Inches(4.9), CW, Inches(0.8),
        "Adopté. Découpage chronologique plus rapide ET qui préserve l'ordre du récit. Base de tout le pipeline.")
analyse(s, CL, Inches(5.75), CW, Inches(1.3),
        "HDBSCAN regroupe par densité sans tenir compte du temps (clusters non contigus, points "
        "isolés jetés comme bruit). La détection de frontières garde l'ordre linéaire et ne coupe "
        "qu'où la cohésion chute : c'est le principe de TextTiling (Hearst 1997), ici sur embeddings.", size=11)
sources(s, CL, Inches(7.05), CW, [
    ("TextTiling, Hearst 1997", "https://aclanthology.org/J97-1003/"),
    ("HDBSCAN, Campello et al. 2013", "https://link.springer.com/chapter/10.1007/978-3-642-37456-2_14"),
])

# 1.2 params
s = content_slide("1.2  Paramètres et résultats réels du découpage retenu")
table(s, CL, Inches(1.15), Inches(9.2),
      ["Paramètre", "Valeur", "Rôle"],
      [["Fenêtre glissante", "3 segments (slide 1)", "granularité de l'analyse"],
       ["Lissage gaussien sigma", "2,0", "atténue le bruit des similarités"],
       ["Percentile de vallée", "5 %", "seuil de détection des frontières"],
       ["Distance min. frontières", "10 fenêtres", "évite les coupures trop proches"],
       ["Taille max. d'un chunk", "15 000 caractères", "déclenche le re-split récursif"],
       ["Embeddings", "all-MiniLM-L6-v2 (384d, CPU)", "vectorisation des fenêtres"]],
      col_fracs=[2.3, 2.4, 3], font=10.5)
analyse(s, CL, Inches(4.35), CW, Inches(1.4),
        "all-MiniLM-L6-v2 est un MiniLM (Wang et al. 2020, distillation d'attention) affiné façon "
        "Sentence-BERT (Reimers et Gurevych 2019) pour que la similarité cosine reflète le sens. "
        "Sur dicte_audio_3 : 264 fenêtres, 6 frontières, 10 chunks ; coût 12 s d'embeddings (négligeable).", size=11)
sources(s, CL, Inches(5.8), CW, [
    ("Sentence-BERT, EMNLP 2019", "https://arxiv.org/abs/1908.10084"),
    ("MiniLM, NeurIPS 2020", "https://arxiv.org/abs/2002.10957"),
])

# 1.3 sans découpage
s = content_slide("1.3  Approches sans découpage (contre-exemples)")
table(s, CL, Inches(1.2), CW,
      ["Pipeline", "Principe", "Modèle(s)", "Durée", "Remarque"],
      [["Multi-pass V1", "4 extractions/chunk + 3 sections (~40 appels)", "Mistral 7B", "1 h 46", "lent"],
       ["Multi-pass V1", "idem", "Qwen2.5 3B", "43 min 38", ""],
       ["Multi-pass V1", "idem", "LFM2.5-1.2B", "18 min 14", "rapport quasi vide"],
       ["3-calls", "0 chunking, 3 appels (ctx 32k)", "LFM Extract+Transcript", "6 min 10", "le plus rapide"],
       ["one-shot", "tout en 1 appel", "LFM2-2.6B", "29 min 09", ""],
       ["one-shot", "tout en 1 appel", "LFM2.5-1.2B", "4 min 17", "1 509 car. (court)"]],
      col_fracs=[1.4, 3, 2, 1.2, 1.6], font=10)
verdict(s, CL, Inches(5.3), CW, Inches(1.2),
        "Les approches sans découpage sont spectaculairement rapides mais produisent des CR trop "
        "courts ou incomplets sur 1 h. Le découpage par frontières reste le meilleur compromis "
        "exhaustivité / fiabilité.")

# ==========================================================================
# PARTIE 2
# ==========================================================================
divider("Partie 2", "Les architectures",
        "Extraction (assemblage déterministe) puis refonte agentique")

# 2.A extraction
s = content_slide("2.A  Architecture extraction + assemblage déterministe")
paragraph(s, CL, Inches(1.1), CW,
          "Par chunk : résumé JSON + extraction JSON (few-shot négatifs), puis assemblage 100 % "
          "Python. Mesuré sur dicte_audio_3 (266 seg, 11 chunks).", size=13)
table(s, CL, Inches(1.85), Inches(11),
      ["Run", "Modèle", "Durée", "RAM pic", "Décisions", "Actions"],
      [["boundary (V3)", "Ministral 3B", "29 min 15", "2 945 Mo", "n.c.", "n.c."],
       ["extraction (V4)", "Ministral 3B", "51 min 52", "5 825 Mo", "0", "0"],
       ["extraction (V4)", "Qwen2.5 3B", "16 min 10", "4 165 Mo", "5", "9"]],
      col_fracs=[1.6, 1.6, 1.4, 1.3, 1.2, 1.2], font=10.5,
      green={(2, 2)}, red={(1, 2)})
verdict(s, CL, Inches(3.7), CW, Inches(0.8),
        "Adopté (extraction JSON + assemblage déterministe), pour la fiabilité du format. "
        "Qwen2.5 3B est environ 3 fois plus rapide que Ministral 3B et ose extraire des items.")
note(s, CL, Inches(4.6), CW, Inches(1.5),
     "Attention : quantité n'est pas qualité. Les 5 décisions de Qwen sont en réalité 2 récits d'un "
     "projet passé de Maya, 1 non-décision et 2 suggestions. Le 0 de Ministral est plus FIDÈLE "
     "(c'est une prise de contact sans décision). Extraire beaucoup peut vouloir dire halluciner beaucoup.",
     size=11.5)

# 2.A side by side V4
s = content_slide("2.A  Même passage : Qwen invente des décisions, Ministral reste fidèle")
paragraph(s, CL, Inches(1.1), CW,
          "L'outil PPT de Maya est un projet PASSÉ qu'elle raconte. Qwen le transforme en décisions "
          "de la réunion ; Ministral indique correctement qu'aucune décision n'a été prise.", size=12)
two_col(s, Inches(1.85),
        "Qwen2.5 3B : fabrique des décisions",
        "« Décisions : Automatisation de la chaîne de traitement... ; Définition d'une structure "
        "préétablie pour le rapport. » Or ce sont des éléments d'un projet ANTÉRIEUR, pas des "
        "décisions prises en séance.",
        "archi_extraction_v4_qwen2.5-3b.md",
        "Ministral 3B : reste fidèle",
        "« ... Aucune décision prise concernant la validation ou l'extension de ces éléments. » "
        "(synthèse : « Si aucune décision claire n'a été prise pour orienter une mission... »)",
        "archi_extraction_v4_ministral3b.md",
        sentA="mauvais", sentB="bon", height=Inches(3.4))

# 2.A models
s = content_slide("2.A  Modèles évalués (≤ 7B, CPU)")
table(s, CL, Inches(1.2), Inches(10),
      ["Modèle", "Famille", "Quantization", "Taille"],
      [["Ministral-3-3B-Instruct-2512", "Mistral 3B", "Q4_K_M", "~1,9 Go"],
       ["Mistral-7B-Instruct-v0.3", "Mistral 7B", "Q4_K_M", "~4,1 Go"],
       ["qwen2.5-3b-instruct", "Qwen 2.5 3B", "Q4_0", "~1,9 Go"],
       ["Qwen3-4B (thinking)", "Qwen 3 4B", "Q4_K_M", "~2,5 Go"],
       ["SmolLM3-3B", "HuggingFace", "Q4_K_M", "~1,9 Go"],
       ["LFM2.5-1.2B et LFM2-2.6B", "Liquid", "Q4 / Q6", "0,66 à 2,0 Go"],
       ["NuExtract-2.0-2B", "Extraction structurée", "Q4_K_M", "~1,3 Go"]],
      col_fracs=[3, 2.2, 1.6, 1.4], font=11)

# 2.B agentique agents
s = content_slide("2.B  Refonte agentique : orchestrateur + workers")
paragraph(s, CL, Inches(1.1), CW,
          "Limite de l'extraction : même structure de CR pour toutes les réunions. L'agentique "
          "adapte la structure au type de réunion via des agents spécialisés.", size=13)
bullets(s, CL, Inches(1.9), CW, Inches(3.0), [
    ("Agent CONTEXT BUILDER : ", "type de réunion (multi-hypothèses), objectif, synthèse globale."),
    ("Agent PLANNER : ", "conçoit la structure (les sections)."),
    ("Agent CONTENT DESIGNER : ", "pour chaque section, choisit les chunks et rédige le brief."),
    ("Agents WORKERS : ", "rédigent chaque section (narratif, puces, tableaux) à partir du texte brut."),
    ("Juges déterministes : ", "décisions et actions, puis assemblage Markdown 100 % Python."),
], size=14)

# 2.B V1-V10 table
s = content_slide("2.B  Évolution version par version (V1 à V10)")
table(s, CL, Inches(1.15), CW,
      ["V", "Modèle(s)", "Sect.", "Durée", "Résultat clé, ce qu'elle corrigeait"],
      [["V1", "Ministral 3B", "3", "20 min", "1er jet : orchestrateur et workers"],
       ["V2", "Ministral 3B", "3", "26 min", "multi-agents ; bug texte brut"],
       ["V3", "Ministral 3B", "4", "76 min", "texte brut OK ; tour de table mal placé"],
       ["V4", "Ministral 3B", "4", "69 min", "tables décisions/actions ; speaker mapping"],
       ["V5", "Qwen7B+draft", "10", "incomplet", "test spec decoding, run interrompu"],
       ["V6", "Qwen 7B", "8", "146 min", "prose en hausse mais détail PPT perdu"],
       ["V7", "Ministral 3B", "5", "182 min", "worker planté (JSON vide), d'où robustesse"],
       ["V8", "Ministral 3B", "6", "125 min", "robustesse + juges few-shot ; extraction fraîche"],
       ["V9", "Hybride 7B+3B", "7", "112 min", "routage par agent (contexte 7B, reste 3B)"],
       ["V10", "Qwen 3B sans diar.", "10", "134 min", "tout-Qwen, transcript sans diarisation"]],
      col_fracs=[0.5, 1.7, 0.6, 1.2, 4.2], font=9.5, row_h=0.31)

# 2.B side by side V8/V6
s = content_slide("2.B  Côte à côte : Ministral 3B (V8) vs Qwen 7B (V6)")
paragraph(s, CL, Inches(1.1), CW,
          "Même section. Le 7B cadre mieux globalement mais perd le détail concret ; le 3B retient "
          "l'outil Power et les fichiers PPT.", size=12)
two_col(s, Inches(1.75),
        "Ministral 3B (V8) : détail conservé",
        "« ... automatiser la rédaction de rapports métiers... L'équipe utilisait déjà un outil "
        "interne nommé Power... résultats stockés en fichiers PPT contenant principalement des "
        "graphiques et peu de texte... »",
        "agentique_v8_ministral3b.md",
        "Qwen 7B (V6) : générique, détail perdu",
        "« ... présentation des rôles et expertises... Ensuite Maya SAHRAOUI a présenté ses "
        "méthodes pour automatiser la rédaction de rapports à partir d'études PowerPoint. »",
        "agentique_v6_qwen7b.md",
        sentA="bon", sentB="mauvais", height=Inches(3.3))

# 2.B defects
s = content_slide("2.B  Défauts (même réunion) : pourquoi ces versions sont écartées")
paragraph(s, CL, Inches(1.05), CW,
          "V10 (tout-Qwen, sans diarisation) : décisions fabriquées + fuite d'un exemple du prompt.", size=12, bold=True)
excerpt(s, CL, Inches(1.55), CW, Inches(1.75),
        ["| 1 | on décide de faire des connaissances | On décide de faire des connaissances |",
         "| 2 | on décide de gagner en souveraineté  | On décide de gagner en souveraineté |",
         "Plan d'action : | 1 | S'occuper de la doc, vendredi | Alice | vendredi |  (Alice vient du prompt)"],
        "agentique_v10_qwen3b_sans-diarisation.md", size=10)
paragraph(s, CL, Inches(3.5), CW,
          "V7 (Ministral 3B) : un worker renvoie un JSON vide, section perdue. Déclencheur des "
          "garde-fous (timeout, retry, fallback déterministe).", size=12, bold=True)
excerpt(s, CL, Inches(4.0), CW, Inches(1.1),
        ["## 5. Plan d'action",
         "_Section non rendue (Expecting value: line 1 column 1 (char 0))._"],
        "agentique_v7_ministral3b.md", size=10)

# 2.B V9 routing
s = content_slide("2.B  V9 : routage hybride par agent")
paragraph(s, CL, Inches(1.1), CW,
          "Le 7B n'est meilleur que sur le cadrage global (Context Builder) et la fluidité ; le 3B "
          "suffit et est plus rapide ailleurs. V9 met un modèle différent par phase.", size=13)
table(s, CL, Inches(1.95), Inches(8),
      ["Phase", "Modèle", "Durée"],
      [["Context Builder", "Qwen 7B", "837 s"],
       ["Planner", "Ministral 3B", "629 s"],
       ["Content Designers", "Ministral 3B", "3 147 s"],
       ["Workers", "Ministral 3B", "2 040 s"],
       ["Swaps de modèle (coût)", "n.c.", "8,1 s"],
       ["Total", "n.c.", "112 min (6 701 s)"]],
      col_fracs=[3, 2, 2], font=11)
verdict(s, CL, Inches(5.3), CW, Inches(1.0),
        "Prometteur (travail en cours). Le routage fonctionne (swap environ 8 s, négligeable) : "
        "qualité du 7B sur le cadrage et vitesse du 3B sur le reste.")

# ==========================================================================
# PARTIE 3
# ==========================================================================
divider("Partie 3", "Les optimisations",
        "Tout ce qui a été testé pour la vitesse et la qualité")

# 3.1 + 3.2
s = content_slide("3.1 et 3.2  Baseline, puis plan d'action perchunk")
paragraph(s, CL, Inches(1.1), Inches(6),
          "Baseline (prompts verbeux) : 64 min 44, 8 appels LLM. Objectif des expériences : "
          "réduire ce temps sans perdre en qualité.", size=12)
table(s, CL, Inches(2.0), CW,
      ["Mode (plan d'action)", "Durée", "Appels LLM", "Items de plan extraits"],
      [["legacy (1 appel final)", "35 min 47", "16 (2/chunk)", "4"],
       ["perchunk (par chunk)", "40 min 33", "29 (4/chunk)", "17 (9 engagements + 8 suggestions)"]],
      col_fracs=[2.2, 1.6, 1.6, 3], font=11,
      green={(1, 3)}, red={(0, 3)})
verdict(s, CL, Inches(3.7), CW, Inches(1.0),
        "Adopté (perchunk). Plus lent d'environ 5 min mais extrait 17 items au lieu de 4, et "
        "l'assemblage final ne coûte plus d'appel LLM. Répond au reproche pas d'action prise.")

# 3.3 flash attention
s = content_slide("3.3  Flash Attention : ON vs OFF")
table(s, CL, Inches(1.15), CW,
      ["Configuration", "Durée", "RAM serveur pic", "Chunks", "Appels"],
      [["Flash Attention ON", "25 min 14 (1 513,8 s)", "4 557 Mo", "10", "31"],
       ["Flash Attention OFF", "33 min 29 (2 009,2 s)", "5 467 Mo", "7", "22"]],
      col_fracs=[2.2, 2.4, 1.8, 1, 1], font=11,
      green={(0, 1)}, red={(1, 1)})
verdict(s, CL, Inches(2.7), CW, Inches(0.7),
        "Adopté (ON). Plus rapide (environ -25 %) et moins de RAM pic.")
analyse(s, CL, Inches(3.5), CW, Inches(1.7),
        "L'attention standard matérialise une matrice N×N qu'il faut écrire puis relire en mémoire "
        "(coûteux en bande passante, le vrai goulot). FlashAttention (Dao et al. 2022) la calcule "
        "par TUILES sans jamais la matérialiser (IO-aware) : moins d'accès mémoire, moins de RAM. "
        "Dans llama.cpp, l'activer fournit ce noyau fusionné ET débloque la quantification du cache KV.",
        size=11.5)
sources(s, CL, Inches(5.3), CW, [
    ("FlashAttention, Dao et al., NeurIPS 2022", "https://arxiv.org/abs/2205.14135"),
])

# 3.4 KV + 3.5 document-first table
s = content_slide("3.4  Cache KV quantifié (q8_0)")
verdict(s, CL, Inches(1.15), CW, Inches(0.7),
        "Adopté. Allège la mémoire du cache sans dégradation observée de la sortie.")
analyse(s, CL, Inches(2.0), CW, Inches(1.9),
        "Pendant la génération, le cache KV grossit avec le contexte et finit par dominer la "
        "mémoire ; or sur CPU le facteur limitant est la bande passante mémoire (relire ce cache à "
        "chaque token). Le quantifier en 8 bits (q8_0) divise environ par deux son empreinte et la "
        "bande passante nécessaire, avec une perte négligeable à 8 bits. C'est le constat général "
        "de la littérature sur la quantification du cache KV (KVQuant).", size=12)
sources(s, CL, Inches(4.0), CW, [
    ("KVQuant, Hooper et al., NeurIPS 2024", "https://arxiv.org/abs/2401.18079"),
])

# 3.5 document-first
s = content_slide("3.5  Ordre du prompt document-first (réutilisation du cache KV)")
paragraph(s, CL, Inches(1.1), CW,
          "Placer le texte du chunk AVANT les instructions. Son cache KV est calculé une fois puis "
          "réutilisé par les appels suivants du même chunk (preuve : chunk 9, run ik_v5).", size=12)
table(s, CL, Inches(1.95), CW,
      ["Appel sur le chunk", "Tokens à prefiller", "Temps de prefill", "Pourquoi"],
      [["1er, Résumé", "1 394 tokens", "~45,5 s", "charge tout le texte du chunk"],
       ["2e, Extraction", "220 tokens", "~9,3 s", "texte déjà en cache, traite juste les instructions"]],
      col_fracs=[2, 1.8, 1.6, 3.2], font=11,
      red={(0, 2)}, green={(1, 2)})
analyse(s, CL, Inches(3.6), CW, Inches(1.6),
        "L'attention est causale : le cache KV d'un token ne dépend que de ce qui le précède. Le "
        "document en tête devient un PRÉFIXE stable, calculé une fois puis réutilisé ; seul le "
        "suffixe d'instructions change. C'est la réutilisation par préfixe formalisée par Prompt "
        "Cache (Gim et al. 2024). Environ 6 fois moins de tokens au 2e appel ; qualité identique.", size=11.5)
sources(s, CL, Inches(5.3), CW, [
    ("Prompt Cache, Gim et al., MLSys 2024", "https://arxiv.org/abs/2311.04934"),
])

# 3.6 ik_llama
s = content_slide("3.6  Backend ik_llama.cpp vs llama.cpp mainline")
table(s, CL, Inches(1.15), CW,
      ["Backend", "Durée", "RAM serveur (moy / pic)", "Décisions/Actions"],
      [["mainline et FA", "25 min 14 (1 513,8 s)", "4 250 / 4 557 Mo", "0 / 6"],
       ["ik_llama v5", "22 min 41 (1 362,5 s)", "2 597 / 3 128 Mo", "3 / 9"]],
      col_fracs=[1.8, 2.4, 2.4, 1.8], font=11,
      green={(1, 1), (1, 2)})
verdict(s, CL, Inches(2.7), CW, Inches(0.8),
        "Prometteur. Plus rapide (environ -10 %) et surtout -1,4 Go de RAM pic. Mainline gardé en "
        "production pour la compatibilité.")
analyse(s, CL, Inches(3.6), CW, Inches(1.6),
        "Fork de llama.cpp (par I. Kawrakow). Apporte des quantifications SOTA (familles IQ et "
        "K-quants améliorées, meilleure perplexité à bits égaux), des noyaux CPU SIMD plus rapides "
        "et un row-interleaved packing des poids (meilleure efficacité du cache processeur). Gains "
        "surtout visibles en CPU. Pas de publication académique : on cite le dépôt.", size=11.5)
sources(s, CL, Inches(5.3), CW, [
    ("ik_llama.cpp, dépôt GitHub (ikawrakow)", "https://github.com/ikawrakow/ik_llama.cpp"),
])

# 3.7 spec decoding
s = content_slide("3.7  Speculative decoding (test dédié)")
table(s, CL, Inches(1.15), Inches(10),
      ["Configuration", "Décodage", "Temps mur", "Acceptance"],
      [["A, Ministral 3B seul", "4,82 tok/s", "268 s", "n.c."],
       ["B, Qwen 7B seul", "3,75 tok/s", "426 s", "n.c."],
       ["C, Qwen 7B + draft 0,5B", "2,00 tok/s", "472 s", "0,64 (61/95)"]],
      col_fracs=[3, 1.6, 1.4, 1.6], font=11,
      green={(0, 1)}, red={(2, 1), (2, 2)})
verdict(s, CL, Inches(3.0), CW, Inches(0.7),
        "Rejeté dans notre configuration (Qwen 7B + draft 0,5B sur CPU). Config C la PLUS LENTE "
        "(2,00 vs 3,75 tok/s).")
analyse(s, CL, Inches(3.85), CW, Inches(2.0),
        "Le speculative decoding (Leviathan et al. 2023) accélère un décodage MEMORY-BOUND : sur "
        "GPU à lot=1 il reste du calcul libre pour vérifier les K tokens du draft presque "
        "gratuitement (2 à 3 fois). Sur CPU ça ne paie pas : (1) la vérification par lots est "
        "limitée par le calcul (peu d'unités parallèles) ; (2) le draft occupe les mêmes cœurs ; "
        "(3) acceptance 0,64 insuffisante. Retours externes mitigés ; notre mesure : ralentissement.",
        size=11)
sources(s, CL, Inches(5.95), CW, [
    ("Speculative Decoding, Leviathan et al., ICML 2023", "https://arxiv.org/abs/2211.17192"),
])

# 3.8 PLD
s = content_slide("3.8  Prompt Lookup Decoding (PLD)")
table(s, CL, Inches(1.15), Inches(10),
      ["Étape", "PLD ON", "PLD OFF", "Écart"],
      [["Extraction et plan", "210,20 s", "209,79 s", "+0,4 s (ON plus lent)"],
       ["Résumé", "224,70 s", "214,02 s", "+10,7 s (ON plus lent)"]],
      col_fracs=[2, 1.6, 1.6, 2.6], font=11)
verdict(s, CL, Inches(2.5), CW, Inches(0.7),
        "Rejeté. Aucun gain : la sortie (résumé abstractif, JSON court) ne recopie pas assez le prompt.")
analyse(s, CL, Inches(3.35), CW, Inches(1.7),
        "Le PLD remplace le draft par une recopie de n-grammes déjà présents dans le prompt : il "
        "n'accélère que si la sortie REPREND littéralement l'entrée (résumé extractif, QA, code). "
        "Notre sortie est ABSTRACTIVE (reformulée) plus un JSON court : recouvrement de n-grammes "
        "faible, presque rien à deviner par copie, donc pas de gain.", size=11.5)
sources(s, CL, Inches(5.15), CW, [
    ("Prompt Lookup Decoding, dépôt GitHub (apoorvumang)", "https://github.com/apoorvumang/prompt-lookup-decoding"),
])

# 3.8 LLMLingua + parallel/openvino
s = content_slide("3.8  Compression LLMLingua-2 ; parallélisme ; OpenVINO")
table(s, CL, Inches(1.15), Inches(11),
      ["Ratio LLMLingua", "Caractères", "Durée totale", "Effet qualité"],
      [["1.0 (aucune)", "8 000", "308,7 s", "référence, extraction vide"],
       ["0.5", "4 152", "174,7 s", "identique à la référence"],
       ["0.4", "3 351", "128,7 s", "1 décision apparaît (extrapolation ?)"]],
      col_fracs=[2, 1.6, 1.6, 3.5], font=11)
verdict(s, CL, Inches(2.9), CW, Inches(0.9),
        "Écarté. Gain net (jusqu'à environ 2,9 fois) mais aux ratios agressifs le contenu diverge "
        "(items non vérifiés). Risque qualité non maîtrisé sans gold.")
bullets(s, CL, Inches(4.0), CW, Inches(1.3), [
    ("Parallélisme (--parallel) : ", "pas de gain réel sur CPU (cœurs partagés). Pipeline séquentiel."),
    ("OpenVINO : ", "non concluant (prérequis lourds, Q4_K_M partiellement supporté)."),
], size=12.5)
sources(s, CL, Inches(5.4), CW, [
    ("LLMLingua-2, Pan et al., Findings of ACL 2024", "https://aclanthology.org/2024.findings-acl.57/"),
])

# ==========================================================================
# PARTIE 4
# ==========================================================================
divider("Partie 4", "Synthèse", "Bilan, choix de modèle, état actuel")

# 4.1 bilan
s = content_slide("4.1  Bilan technique par technique")
table(s, CL, Inches(1.1), CW,
      ["Technique", "Partie", "Verdict"],
      [["Découpage par frontières (vs HDBSCAN)", "1.1", "ADOPTÉ, chronologique, -32 à -41 %"],
       ["Extraction JSON + assemblage déterministe", "2.A", "ADOPTÉ, fiabilité du format"],
       ["Plan d'action perchunk", "3.2", "ADOPTÉ, 17 items vs 4"],
       ["Flash Attention", "3.3", "ADOPTÉ, environ -25 %, moins de RAM"],
       ["Cache KV quantifié (q8_0)", "3.4", "ADOPTÉ, bande passante CPU"],
       ["Prompt document-first (cache KV)", "3.5", "ADOPTÉ, prefill 6 fois moindre au 2e appel"],
       ["Backend ik_llama.cpp", "3.6", "PROMETTEUR, plus rapide, -1,4 Go"],
       ["Speculative decoding", "3.7", "REJETÉ sur CPU, ralentit"],
       ["Prompt Lookup Decoding", "3.8", "REJETÉ, aucun gain"],
       ["Compression LLMLingua", "3.8", "ÉCARTÉ, risque qualité"],
       ["Parallélisme / OpenVINO", "3.8", "REJETÉ / NON CONCLU"],
       ["Architecture agentique + routage hybride", "2.B", "EN COURS, V8/V9"]],
      col_fracs=[3.4, 0.8, 3.4], font=10, row_h=0.3)

# 4.2 choix modèle
s = content_slide("4.2  Choix de modèle : ce que montrent les benchmarks")
bullets(s, CL, Inches(1.3), CW, Inches(4.5), [
    ("Qwen2.5 3B : ", "le plus rapide des 3B et ose extraire, mais sur-extrait (fabrique des items quand il n'y en a pas). À encadrer par un juge."),
    ("Ministral 3B : ", "modèle français de référence, fiable et sobre (retenue = atout anti-hallucination), plus faible sur le cadrage global."),
    ("Qwen 7B : ", "meilleur cadrage et fluidité, mais environ 2 fois plus lent ET perd du détail concret. Réservé au Context Builder."),
    ("Modèles < 1,5B (LFM2.5, SmolLM) : ", "très rapides mais comptes rendus trop courts et pauvres."),
], size=15)

# 4.3 état
s = content_slide("4.3  État actuel et pistes")
box(s, CL, Inches(1.3), CW, Inches(1.5), "Production",
    "Boundary + extraction JSON + assemblage déterministe, Ministral 3B Q4_K_M, Flash Attention ON, "
    "cache KV q8_0, plan perchunk, document-first, séquentiel. Environ 25 min sur 1 h de réunion, "
    "moins de 5 Go de RAM, 100 % local.", VERT_BG, VERT, size=13)
box(s, CL, Inches(3.0), CW, Inches(1.4), "En cours",
    "Refonte agentique (V8/V9) pour adapter la structure au type de réunion ; routage hybride "
    "(7B cadrage, 3B reste) ; garde-fous de robustesse (timeout, retry, fallback, cap sur le nombre de tableaux).",
    BLEU_BG, BLEU, size=13)
box(s, CL, Inches(4.6), CW, Inches(1.5), "Pistes",
    "Backend ik_llama (RAM) ; amélioration ASR et diarisation en amont (cause racine de plusieurs "
    "erreurs de noms) ; si la contrainte 100 % local était levée, variante cloud map-reduce "
    "(extraction locale + synthèse cloud) plus rapide et de meilleure qualité.", JAUNE_BG, RGBColor(0x8A, 0x6D, 0x00), size=13)

# Sources récap
s = content_slide("Sources académiques (vérifiées)")
items = [
    ("TextTiling, Hearst, Computational Linguistics 1997", "https://aclanthology.org/J97-1003/"),
    ("HDBSCAN, Campello, Moulavi, Sander, PAKDD 2013", "https://link.springer.com/chapter/10.1007/978-3-642-37456-2_14"),
    ("Sentence-BERT, Reimers et Gurevych, EMNLP 2019", "https://arxiv.org/abs/1908.10084"),
    ("MiniLM, Wang et al., NeurIPS 2020", "https://arxiv.org/abs/2002.10957"),
    ("FlashAttention, Dao et al., NeurIPS 2022", "https://arxiv.org/abs/2205.14135"),
    ("KVQuant, Hooper et al., NeurIPS 2024", "https://arxiv.org/abs/2401.18079"),
    ("Prompt Cache, Gim et al., MLSys 2024", "https://arxiv.org/abs/2311.04934"),
    ("Speculative Decoding, Leviathan et al., ICML 2023", "https://arxiv.org/abs/2211.17192"),
    ("Prompt Lookup Decoding, GitHub apoorvumang", "https://github.com/apoorvumang/prompt-lookup-decoding"),
    ("LLMLingua-2, Pan et al., Findings of ACL 2024", "https://aclanthology.org/2024.findings-acl.57/"),
    ("ik_llama.cpp, GitHub ikawrakow", "https://github.com/ikawrakow/ik_llama.cpp"),
]
tf = textbox(s, CL, Inches(1.2), CW, Inches(5.8))
for i, (label, url) in enumerate(items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(7)
    rb = p.add_run()
    rb.text = "•  "
    rb.font.size = Pt(13)
    rb.font.color.rgb = BLEU2
    r = p.add_run()
    r.text = label
    r.font.size = Pt(13)
    r.font.color.rgb = BLEU
    r.font.underline = True
    r.hyperlink.address = url

# ==========================================================================
# Sauvegarde
# ==========================================================================
out = Path(__file__).parent / "Presentation_Partie_LLM.pptx"
try:
    prs.save(out)
    print(f"Presentation generee : {out}  ({len(prs.slides._sldIdLst)} slides)")
except PermissionError:
    alt = Path(__file__).parent / "Presentation_Partie_LLM_NEW.pptx"
    prs.save(alt)
    print(f"[!] {out.name} verrouille -> enregistre sous : {alt}")
