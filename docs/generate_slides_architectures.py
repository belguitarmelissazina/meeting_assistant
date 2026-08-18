"""Genere un PPTX comparant les 2 architectures de generation de CR :

Architecture 1 — Mainline (en dur, structure fixe)
Architecture 2 — Agentique (orchestrateur + workers + renderer)

Pour chacune, on detaille :
- ce qui est parallelisable avec l'enregistrement et ce qui ne l'est pas
- chaque appel LLM ("agent") : ce qu'il fait, son prompt, son JSON in/out
- comment se fait l'assemblage
- le temps en minutes

Sources des donnees :
- Architecture 1 : _bench_prod_mainline_fa/compte_rendu.metrics.json
- Architecture 2 : _bench_prod_mainline_fa/orchestrator_v1.json (timings)
  + on additionne car l'agentique reutilise les sections deja generees
"""

from __future__ import annotations

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor


# ============================================================================
# Constantes (donnees benchmarks)
# ============================================================================

# --- Architecture 1 — Mainline FA (bench du 2026-05-13) -------------------
ARCHI1 = {
    "label": "Architecture 1 — Mainline (structure fixe, en dur)",
    "model": "Ministral 3B Q4_K_M (llama.cpp, flash attention)",
    "n_chunks": 10,
    "n_llm_calls": 31,
    "timings_min": {
        "embeddings": 7.4 / 60,
        "construction_chunks": 5.5 / 60,
        "generation_sections_llm": 1233.5 / 60,      # 3 appels par chunk
        "executive_summary_llm": 257.6 / 60,
        "assemblage_deterministe": 0.003 / 60,
    },
    "total_min": 1513.8 / 60,                          # 25.23 min
    "parallelisable_min": 1233.5 / 60,                 # pendant l'enregistrement
    "sequentiel_min": (257.6 + 0.003) / 60,            # apres l'enregistrement
}

# --- Architecture 2 — Agentique (bench orchestrator_v1.json) --------------
# Total agentique = mainline (extraction des sections) + orchestrateur + workers
ARCHI2_ADD_MIN = 755.97 / 60                            # = 12.60 min
ARCHI2 = {
    "label": "Architecture 2 — Agentique (orchestrateur + workers)",
    "model": "Ministral 3B Q4_K_M (llama.cpp, flash attention)",
    "n_chunks": 10,
    "n_sections_plan": 4,
    "n_llm_calls_agentique": 1 + 4,                    # 1 orch + 4 workers
    "n_llm_calls_total": 31 + 5,                       # extraction + agentique
    "timings_min": {
        "llama_server_startup": 12.24 / 60,
        "orchestrateur": 525.22 / 60,
        "workers (sequentiels)": 218.49 / 60,
        "render markdown (python)": 0.0 / 60,
    },
    "total_min": ARCHI1["total_min"] + ARCHI2_ADD_MIN,
    "parallelisable_min": ARCHI1["parallelisable_min"],
    "sequentiel_min": ARCHI1["sequentiel_min"] + ARCHI2_ADD_MIN,
    # Note : la user veut le total = bench mainline FA + bench orchestrateur
    # car les sections.json sont reutilisees.
}


# --- Prompts mainline (extraits de compte_rendu.metrics.json) -------------
PROMPT_RESUME = """Extrait :
{texte}

---

Resume cet extrait de reunion en UN objet JSON.

REGLES :
- titre : court et precis, sujet principal de l'extrait
- contexte : UNE phrase qui resume le sujet
- points : 3 a 5 bullets factuels, chacun UNE phrase courte autonome (qui,
  quoi, pourquoi). Pas de transition ("ensuite", "puis"), pas de redite du
  contexte.
- Style sobre, pas d'opinion.
- N'INVENTE RIEN. Ne developpe pas les sigles.

FORMAT JSON STRICT, rien d'autre :
{"titre": "string", "contexte": "string", "points": ["string", ...]}"""

PROMPT_EXTRACTION = """Extrait :
{texte}

---

Extrais les decisions de cet extrait de reunion. Produis UN objet JSON.

REGLE ABSOLUE : EXTRACTION PURE. Ne liste QUE ce qui est EXPLICITEMENT dit.

DECISION = choix ACTE par le groupe ("on decide de", "on valide", "c'est acte").
- Une discussion ou exploration d'options N'EST PAS une decision.
- Un engagement a faire une action future N'EST PAS une decision.
- Si rien n'est acte, tableau vide.

EXEMPLES :
"J'ai developpe un outil N8N avec plusieurs agents." -> {"decisions": []}
"On valide le passage en production la semaine prochaine."
   -> {"decisions": ["Passage en production valide pour la semaine prochaine"]}

FORMAT JSON STRICT, rien d'autre :
{"decisions": ["string", ...]}"""

PROMPT_PLAN_PERCHUNK = """Extrait :
{texte}

---

Extrais les ACTIONS POST-REUNION concretes de cet extrait. JSON.

DEUX TYPES :
1. ENGAGEMENT : action promise EXPLICITEMENT ("je vais X", "d'ici vendredi").
   -> Responsable et echeance EXACTEMENT tels que mentionnes.
2. SUGGESTION : action de bon sens qui decoule du sujet, sans avoir ete
   promise. -> Responsable et echeance toujours "—".

REGLE ABSOLUE — RIEN SI RIEN :
- Mieux vaut un tableau vide qu'une action forcee.
- Recap, smalltalk, tour de table -> [].
- Action faite PENDANT la reunion N'EST PAS post-reunion.
- Mention d'un outil/modele DEJA existant N'EST PAS une action future.

LIMITE STRICTE : 0 a 2 items max.

FORMAT JSON STRICT :
{"plan": [{"action": "...", "responsable": "...", "echeance": "...",
           "type": "engagement | suggestion"}, ...]}"""

PROMPT_EXEC_SUMMARY = """Voici le debut du transcript d'une reunion + les
titres/resumes de chaque section. Redige un Executive Summary de 5 phrases.

OBJECTIF : but global, grands themes, conclusion. PAS de resume section
par section, vision de haut niveau.

REGLES :
- N'invente aucune information.
- Ne developpe pas les sigles.
- Pas de listes, pas de titres : un paragraphe.

--- DEBUT TRANSCRIPT ---
{intro}

--- SECTIONS ---
{contenu}"""

# --- Prompts agentique (extraits de _bench_orchestrator.py) ---------------
PROMPT_ORCHESTRATOR = """Tu es un ORCHESTRATEUR de redaction de compte rendu
de reunion. A partir des resumes COMPLETS de chaque passage, CONCOIS le
plan ET REDIGE le prompt complet de chaque worker.

Tu produis :
1. Le TYPE de reunion (revue projet, brainstorm, daily, prise de contact...)
2. Le PLAN (liste de sections). Pour chaque section :
   - section_id (snake_case)
   - section_title
   - section_type (narrative / list_bullets / table_decisions / table_actions)
   - chunk_ids (quels passages sont utiles)
   - prompt_worker : LE PROMPT COMPLET que le worker recevra (regles
     anti-hallucination, role, ton, focus, format)
3. raisonnement : 2-4 phrases sur le pourquoi du plan

REGLES :
- Au moins 3 sections, autant que necessaire.
- N'invente pas une section "Decisions" s'il n'y a pas eu de decisions.
- chunk_ids doivent exister dans le sommaire.
- prompt_worker doit inclure les regles anti-hallucination, francais, pas
  de developpement de sigles, fallback si pas de matiere.
- NE PAS inclure les passages dans prompt_worker (ajoutes ensuite par
  Python).

--- META REUNION ---
{meeting_meta}
--- SOMMAIRE DES PASSAGES (resume COMPLET) ---
{sections_overview}
--- CHUNK_IDS VALIDES ---
{chunk_ids_valides}"""

PROMPT_WORKER = """[prompt redige par l'orchestrateur, libre et adapte
a la section]

--- PASSAGES PERTINENTS (assignes par l'orchestrateur) ---
=== PASSAGE 0 ===
Titre : ...
Contexte : ...
Points :
  - ...
Decisions explicites :
  - ...
Plan (engagements / suggestions) :
  - [engagement] ... | resp: ... | echeance: ...
=== PASSAGE 2 ===
..."""


# ============================================================================
# Helpers PPTX
# ============================================================================

# Couleurs
COLOR_TITLE = RGBColor(0x1F, 0x4E, 0x79)         # bleu fonce
COLOR_ACCENT_1 = RGBColor(0x2E, 0x75, 0xB6)      # bleu archi 1
COLOR_ACCENT_2 = RGBColor(0xC0, 0x50, 0x4D)      # rouge archi 2
COLOR_SUBTLE = RGBColor(0x55, 0x55, 0x55)
COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)
COLOR_BG_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)
COLOR_GREEN = RGBColor(0x2E, 0x7D, 0x32)
COLOR_ORANGE = RGBColor(0xE6, 0x7E, 0x22)


def add_title(slide, text, color=COLOR_TITLE, top_in=0.2, size=24):
    box = slide.shapes.add_textbox(
        Inches(0.4), Inches(top_in), Inches(12.5), Inches(0.7)
    )
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.runs[0].font.bold = True
    p.runs[0].font.size = Pt(size)
    p.runs[0].font.color.rgb = color
    return box


def add_text(slide, text, left, top, width, height, size=14,
             color=COLOR_TEXT, bold=False, mono=False, bg=None):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)

    if bg is not None:
        box.fill.solid()
        box.fill.fore_color.rgb = bg
        box.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    else:
        box.fill.background()
        box.line.fill.background()

    for i, line in enumerate(text.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        run = p.runs[0] if p.runs else None
        if run is None:
            r = p.add_run()
            r.text = ""
            run = r
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        if mono:
            run.font.name = "Consolas"
    return box


def add_section_title(slide, text, left, top, width, color=COLOR_ACCENT_1):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(0.4)
    )
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.runs[0].font.bold = True
    p.runs[0].font.size = Pt(16)
    p.runs[0].font.color.rgb = color
    return box


def add_box_rect(slide, left, top, width, height, fill=None, line=None):
    rect = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    if fill is None:
        rect.fill.background()
    else:
        rect.fill.solid()
        rect.fill.fore_color.rgb = fill
    if line is not None:
        rect.line.color.rgb = line
        rect.line.width = Pt(1.5)
    rect.shadow.inherit = False
    return rect


def add_arrow(slide, x1, y1, x2, y2, color=COLOR_SUBTLE):
    line = slide.shapes.add_connector(
        1, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = color
    line.line.width = Pt(2)
    return line


def new_slide(prs, layout_idx=6):
    return prs.slides.add_slide(prs.slide_layouts[layout_idx])


# ============================================================================
# Slides
# ============================================================================

def make_prs() -> Presentation:
    prs = Presentation()
    # 16:9 — 13.333 x 7.5 inches
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def slide_cover(prs):
    s = new_slide(prs)
    add_box_rect(s, 0, 0, 13.333, 7.5, fill=COLOR_TITLE)

    box = s.shapes.add_textbox(Inches(0.6), Inches(2.0),
                                Inches(12.2), Inches(1.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Generation de comptes rendus de reunion"
    r = p.runs[0]
    r.font.bold = True
    r.font.size = Pt(40)
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    box2 = s.shapes.add_textbox(Inches(0.6), Inches(3.4),
                                  Inches(12.2), Inches(1.0))
    tf2 = box2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = ("Comparaison de 2 architectures : "
              "Mainline (en dur)  vs  Agentique (orchestrateur + workers)")
    r = p.runs[0]
    r.font.size = Pt(22)
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    box3 = s.shapes.add_textbox(Inches(0.6), Inches(5.5),
                                  Inches(12.2), Inches(1.0))
    tf3 = box3.text_frame
    p = tf3.paragraphs[0]
    p.text = ("Modele : Ministral 3B Q4_K_M  |  100 % local  |  "
              "Reunion test : 10 chunks (~17 min audio)")
    r = p.runs[0]
    r.font.size = Pt(16)
    r.font.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)


# ----- Architecture 1 ---------------------------------------------------

def slide_archi1_vue(prs):
    s = new_slide(prs)
    add_title(s, ARCHI1["label"], color=COLOR_ACCENT_1)

    # Bandeau timing
    add_text(s,
             f"Total : {ARCHI1['total_min']:.1f} min   "
             f"|   Appels LLM : {ARCHI1['n_llm_calls']}   "
             f"|   Chunks : {ARCHI1['n_chunks']}   "
             f"|   {ARCHI1['model']}",
             0.4, 1.0, 12.5, 0.5, size=14, color=COLOR_SUBTLE, bold=True)

    # Phase parallelisable
    add_box_rect(s, 0.4, 1.8, 6.2, 5.0,
                 fill=RGBColor(0xE8, 0xF1, 0xFA),
                 line=COLOR_ACCENT_1)
    add_section_title(s, "[PARALLELISABLE pendant l'enregistrement]",
                       0.6, 1.95, 6.0, color=COLOR_GREEN)
    add_text(s,
             f"Duree : {ARCHI1['parallelisable_min']:.1f} min "
             f"(20.6 min sur 17 min d'audio = legerement plus lent que "
             f"le temps reel)",
             0.6, 2.4, 6.0, 0.6, size=12, color=COLOR_SUBTLE)
    add_text(s,
             "Pour CHAQUE chunk (10 chunks), 3 appels LLM successifs :\n\n"
             "  Agent 1 — RESUME\n"
             "     -> titre + contexte + points (JSON)\n\n"
             "  Agent 2 — EXTRACTION decisions\n"
             "     -> decisions[] (JSON)\n\n"
             "  Agent 3 — PLAN per chunk\n"
             "     -> plan[] (engagements + suggestions, JSON)\n\n"
             "Soit 3 x 10 = 30 appels LLM.",
             0.6, 3.0, 6.0, 3.8, size=14)

    # Phase sequentielle
    add_box_rect(s, 6.8, 1.8, 6.2, 5.0,
                 fill=RGBColor(0xFD, 0xEE, 0xE8),
                 line=COLOR_ACCENT_2)
    add_section_title(s, "[SEQUENTIEL apres l'enregistrement]",
                       7.0, 1.95, 6.0, color=COLOR_ORANGE)
    add_text(s,
             f"Duree : {ARCHI1['sequentiel_min']:.1f} min",
             7.0, 2.4, 6.0, 0.4, size=12, color=COLOR_SUBTLE)
    add_text(s,
             "Une fois la reunion finie :\n\n"
             "  Agent 4 — EXECUTIVE SUMMARY\n"
             "     INPUT : intro transcript + tous les resumes des "
             "sections\n"
             "     OUTPUT : 1 paragraphe de 5 phrases\n\n"
             "  Assemblage DETERMINISTE (Python pur, 0 LLM)\n"
             "     - Construit le markdown avec structure FIXE :\n"
             "       Participants / Executive Summary /\n"
             "       Sujets abordes / Decisions / Plan d'attaque\n"
             "     - Plan d'attaque assemble en triant les items\n"
             "       extraits par Agent 3 (engagements puis suggestions)\n\n"
             "Soit 1 appel LLM + 0 ms d'assemblage Python.",
             7.0, 3.0, 6.0, 3.8, size=14)


def slide_archi1_agent(prs, n, label, color, role, prompt, inp, out):
    s = new_slide(prs)
    add_title(s, f"Archi 1 — Agent {n} : {label}", color=color)

    # In / Out (haut)
    add_section_title(s, "Role", 0.4, 0.95, 6, color=COLOR_SUBTLE)
    add_text(s, role, 0.4, 1.30, 12.5, 0.6, size=14)

    add_section_title(s, "INPUT", 0.4, 2.0, 6, color=COLOR_GREEN)
    add_text(s, inp, 0.4, 2.35, 6.2, 1.5, size=11, mono=True,
             bg=RGBColor(0xF5, 0xFA, 0xF5))

    add_section_title(s, "OUTPUT (JSON)", 6.8, 2.0, 6, color=COLOR_ORANGE)
    add_text(s, out, 6.8, 2.35, 6.2, 1.5, size=11, mono=True,
             bg=RGBColor(0xFD, 0xF7, 0xF0))

    add_section_title(s, "PROMPT utilise", 0.4, 3.95, 12, color=COLOR_TITLE)
    add_text(s, prompt, 0.4, 4.30, 12.5, 3.0, size=10, mono=True,
             bg=COLOR_BG_LIGHT)


def slide_archi1_assemblage(prs):
    s = new_slide(prs)
    add_title(s, "Archi 1 — Assemblage mecanique (Python pur, 0 LLM)",
              color=COLOR_ACCENT_1)

    add_text(s,
             "Une fois toutes les sections extraites et l'exec summary "
             "redige, le compte rendu est assemble en Python sans aucun "
             "appel LLM. Structure FIXE imposee, quelque soit le type "
             "de reunion.",
             0.4, 1.0, 12.5, 0.8, size=14)

    add_section_title(s, "Structure FIXE du markdown produit", 0.4, 1.9, 12)
    add_text(s,
             "# Compte rendu de reunion\n\n"
             "## 1. Participants               (si speaker mapping connu)\n"
             "## 2. Executive Summary          (paragraphe Agent 4)\n"
             "## 3. Sujets abordes             (10 sous-sections, "
             "1 par chunk)\n"
             "        ### 1. {titre chunk 0}\n"
             "             {contexte}\n"
             "             - {point 1}\n"
             "             - {point 2}\n"
             "             ...\n"
             "        ### 2. {titre chunk 1}\n"
             "        ...\n"
             "## 4. Decisions                  (concatenation des decisions[])\n"
             "## 5. Plan d'attaque             (tableau GFM, tri engagement "
             ">> suggestion)",
             0.4, 2.3, 12.5, 3.0, size=12, mono=True, bg=COLOR_BG_LIGHT)

    add_section_title(s, "Plan d'attaque — assemblage", 0.4, 5.4, 12,
                       color=COLOR_GREEN)
    add_text(s,
             "build_plan_attack_perchunk() : itere sur les sections, regroupe "
             "tous les items extraits par Agent 3, separe engagements et "
             "suggestions, sort un tableau Markdown avec colonnes "
             "(Action, Sujet, Responsable, Echeance). Duree : 3 ms.",
             0.4, 5.75, 12.5, 1.2, size=12)


# ----- Architecture 2 ---------------------------------------------------

def slide_archi2_vue(prs):
    s = new_slide(prs)
    add_title(s, ARCHI2["label"], color=COLOR_ACCENT_2)

    add_text(s,
             f"Total : {ARCHI2['total_min']:.1f} min   "
             f"(= mainline {ARCHI1['total_min']:.1f} min  "
             f"+  agentique {ARCHI2_ADD_MIN:.1f} min)   "
             f"|   Appels LLM : {ARCHI2['n_llm_calls_total']}   "
             f"(extraction {ARCHI1['n_llm_calls']}  "
             f"+  agentique {ARCHI2['n_llm_calls_agentique']})",
             0.4, 1.0, 12.5, 0.5, size=13, color=COLOR_SUBTLE, bold=True)

    # Phase parallelisable (identique au mainline)
    add_box_rect(s, 0.4, 1.8, 6.2, 5.0,
                 fill=RGBColor(0xE8, 0xF1, 0xFA),
                 line=COLOR_ACCENT_1)
    add_section_title(s, "[PARALLELISABLE pendant l'enregistrement]",
                       0.6, 1.95, 6.0, color=COLOR_GREEN)
    add_text(s,
             f"Duree : {ARCHI2['parallelisable_min']:.1f} min   "
             f"(IDENTIQUE au mainline — on reutilise les sections.json "
             f"deja generees pendant la reunion)",
             0.6, 2.4, 6.0, 0.6, size=12, color=COLOR_SUBTLE)
    add_text(s,
             "Meme phase qu'en mainline :\n\n"
             "  Agent 1 — RESUME       (par chunk)\n"
             "  Agent 2 — EXTRACTION   (par chunk)\n"
             "  Agent 3 — PLAN         (par chunk)\n\n"
             "= 30 appels LLM, produit sections.json :\n"
             "  [\n"
             "    {chunk_id, titre, contexte, points,\n"
             "     resume, decisions, plan},\n"
             "    ...\n"
             "  ]\n\n"
             "C'est l'INPUT de la phase agentique.",
             0.6, 3.0, 6.0, 3.8, size=14)

    # Phase agentique (sequentielle)
    add_box_rect(s, 6.8, 1.8, 6.2, 5.0,
                 fill=RGBColor(0xFD, 0xEE, 0xE8),
                 line=COLOR_ACCENT_2)
    add_section_title(s, "[SEQUENTIEL apres l'enregistrement]",
                       7.0, 1.95, 6.0, color=COLOR_ORANGE)
    add_text(s,
             f"Duree : {ARCHI2_ADD_MIN:.1f} min "
             f"(orchestrateur 8.8 min + workers 3.6 min)",
             7.0, 2.4, 6.0, 0.4, size=12, color=COLOR_SUBTLE)
    add_text(s,
             "Agent ORCHESTRATEUR (1 appel LLM)\n"
             "  -> decide :\n"
             "     - type de reunion\n"
             "     - sections du plan (id, title, type)\n"
             "     - chunk_ids assignes a chaque section\n"
             "     - prompt complet pour chaque worker\n\n"
             "Agent WORKER (1 appel LLM par section)\n"
             "  -> recoit prompt orchestrateur + chunks assignes\n"
             "  -> redige le contenu (JSON typage par section_type)\n"
             "  -> ici 4 sections = 4 appels SEQUENTIELS\n\n"
             "Renderer Python (0 LLM)\n"
             "  -> markdown deterministe a partir des JSON workers\n\n"
             "Total : 1 orch + 4 workers = 5 appels LLM",
             7.0, 3.0, 6.0, 4.0, size=13)


def slide_archi2_orchestrateur(prs):
    s = new_slide(prs)
    add_title(s, "Archi 2 — Agent ORCHESTRATEUR (1 appel LLM)",
              color=COLOR_ACCENT_2)

    add_text(s,
             "Role : decide la structure du compte rendu et redige les "
             "prompts de chaque worker, en UNE SEULE generation LLM.",
             0.4, 1.0, 12.5, 0.6, size=13, color=COLOR_SUBTLE, bold=True)

    add_section_title(s, "INPUT", 0.4, 1.7, 6, color=COLOR_GREEN)
    add_text(s,
             "Pour chaque chunk (10) :\n"
             "  - chunk_id\n"
             "  - titre\n"
             "  - resume COMPLET (non tronque)\n"
             "  - a_decisions (bool)\n"
             "  - nb_engagements (int)\n"
             "  - nb_suggestions (int)\n\n"
             "+ meeting_meta (sujet, participants, duree si dispo)\n"
             "+ chunk_ids_valides",
             0.4, 2.05, 6.2, 1.9, size=11, mono=True,
             bg=RGBColor(0xF5, 0xFA, 0xF5))

    add_section_title(s, "OUTPUT (JSON contraint)", 6.8, 1.7, 6,
                       color=COLOR_ORANGE)
    add_text(s,
             "{\n"
             '  "meeting_type": "prise de contact ...",\n'
             '  "raisonnement": "Ce type combine...",\n'
             '  "plan": [\n'
             "    {\n"
             '      "section_id": "tour_de_table",\n'
             '      "section_title": "Presentation...",\n'
             '      "section_type": "list_bullets",\n'
             '      "chunk_ids": [0],\n'
             '      "prompt_worker": "Role : ..."\n'
             "    },\n"
             "    ... (autres sections)\n"
             "  ]\n"
             "}",
             6.8, 2.05, 6.2, 1.9, size=11, mono=True,
             bg=RGBColor(0xFD, 0xF7, 0xF0))

    add_section_title(s, "PROMPT (extrait)", 0.4, 4.05, 12,
                       color=COLOR_TITLE)
    add_text(s, PROMPT_ORCHESTRATOR, 0.4, 4.40, 12.5, 3.0, size=9, mono=True,
             bg=COLOR_BG_LIGHT)


def slide_archi2_worker(prs):
    s = new_slide(prs)
    add_title(s, "Archi 2 — Agent WORKER (1 appel LLM par section)",
              color=COLOR_ACCENT_2)

    add_text(s,
             "Role : redige le contenu de UNE section en utilisant les "
             "chunks assignes par l'orchestrateur. Sequentiel : section 1 "
             "puis 2 puis 3 puis 4 (pas de parallelisation).",
             0.4, 1.0, 12.5, 0.7, size=13, color=COLOR_SUBTLE, bold=True)

    add_section_title(s, "INPUT", 0.4, 1.8, 6, color=COLOR_GREEN)
    add_text(s,
             "prompt_worker_final =\n"
             "  prompt_worker (redige par l'orchestrateur)\n"
             "  + \"--- PASSAGES PERTINENTS ---\"\n"
             "  + format_chunks_for_worker(chunks_assignes)\n\n"
             "Chunks formates : titre + contexte + points\n"
             "+ decisions explicites + plan items typees\n"
             "(seulement les chunk_ids designes par orch)",
             0.4, 2.15, 6.2, 1.85, size=11, mono=True,
             bg=RGBColor(0xF5, 0xFA, 0xF5))

    add_section_title(s, "OUTPUT (selon section_type)", 6.8, 1.8, 6,
                       color=COLOR_ORANGE)
    add_text(s,
             "narrative      -> {\"content\": \"...\"}\n"
             "list_bullets   -> {\"items\": [\"...\", ...]}\n"
             "table_decisions-> {\"rows\": [\n"
             "                    {\"sujet\":..,\n"
             "                     \"decision\":..}\n"
             "                   ]}\n"
             "table_actions  -> {\"rows\": [\n"
             "                    {\"action\":..,\n"
             "                     \"responsable\":..,\n"
             "                     \"echeance\":..}\n"
             "                   ]}",
             6.8, 2.15, 6.2, 1.85, size=11, mono=True,
             bg=RGBColor(0xFD, 0xF7, 0xF0))

    add_section_title(s, "PROMPT (squelette envoye au worker)",
                       0.4, 4.10, 12, color=COLOR_TITLE)
    add_text(s, PROMPT_WORKER,
             0.4, 4.45, 12.5, 2.85, size=10, mono=True,
             bg=COLOR_BG_LIGHT)


def slide_archi2_renderer(prs):
    s = new_slide(prs)
    add_title(s, "Archi 2 — Renderer Python (0 LLM)",
              color=COLOR_ACCENT_2)

    add_text(s,
             "Une fois les workers termines, l'assemblage du markdown final "
             "est entierement deterministe (pas d'appel LLM). Chaque "
             "section_type a son renderer dedie qui transforme le JSON du "
             "worker en markdown.",
             0.4, 1.0, 12.5, 1.2, size=13, color=COLOR_SUBTLE)

    add_section_title(s, "Mapping section_type -> renderer", 0.4, 2.2, 12,
                       color=COLOR_TITLE)
    add_text(s,
             "narrative       -> \"## N. {title}\\n\\n{content}\\n\"\n\n"
             "list_bullets    -> \"## N. {title}\\n\\n- {item1}\\n"
             "- {item2}\\n...\"\n\n"
             "table_decisions -> \"## N. {title}\\n\\n"
             "| # | Sujet | Decision |\\n|---|---|---|\\n...\"\n\n"
             "table_actions   -> \"## N. {title}\\n\\n"
             "| # | Action | Responsable | Echeance |\\n|---|---|---|---|\\n...\"\n\n"
             "+ numerotation sequentielle\n"
             "+ helper _md_cell() qui escape | et \\n\n"
             "+ etats vides : \"_Aucune decision prise._\", "
             "\"_Aucune action definie._\"",
             0.4, 2.6, 12.5, 3.5, size=13, mono=True, bg=COLOR_BG_LIGHT)

    add_text(s,
             "Duree : ~0 ms. Aucun LLM. Aucune incertitude. "
             "Visuellement comparable au mainline, mais STRUCTURE adaptee "
             "par l'orchestrateur a chaque type de reunion.",
             0.4, 6.3, 12.5, 0.8, size=13, color=COLOR_GREEN, bold=True)


# ----- Comparatif final --------------------------------------------------

def slide_comparatif(prs):
    s = new_slide(prs)
    add_title(s, "Comparatif final", color=COLOR_TITLE)

    rows = [
        ("",                      "Architecture 1",
                                  "Architecture 2"),
        ("",                      "Mainline (en dur)",
                                  "Agentique"),
        ("Structure du CR",       "FIXE (5 sections imposees)",
                                  "ADAPTEE au type de reunion par l'orch"),
        ("Appels LLM total",      "31",
                                  "36 (31 extraction + 1 orch + 4 workers)"),
        ("Phase parallelisable",  f"{ARCHI1['parallelisable_min']:.1f} min "
                                  f"(extraction des 10 chunks pendant la "
                                  f"reunion)",
                                  f"{ARCHI2['parallelisable_min']:.1f} min "
                                  f"(IDENTIQUE — on reutilise sections.json)"),
        ("Phase sequentielle",    f"{ARCHI1['sequentiel_min']:.1f} min "
                                  f"(exec summary 4.3 min + assemblage 0 ms)",
                                  f"{ARCHI2['sequentiel_min']:.1f} min "
                                  f"(exec summary 4.3 min + orch 8.8 min + "
                                  f"workers 3.6 min)"),
        ("Total",                 f"{ARCHI1['total_min']:.1f} min",
                                  f"{ARCHI2['total_min']:.1f} min "
                                  f"(= mainline + agentique)"),
        ("Cout post-reunion",     "4.3 min",
                                  "16.7 min  (~ x4)"),
    ]

    # Tableau natif PPTX
    n_rows = len(rows)
    n_cols = 3
    left = Inches(0.4)
    top = Inches(1.1)
    width = Inches(12.5)
    height = Inches(5.5)
    table = s.shapes.add_table(n_rows, n_cols, left, top, width, height).table

    table.columns[0].width = Inches(3.0)
    table.columns[1].width = Inches(4.5)
    table.columns[2].width = Inches(5.0)

    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = val
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(12)
                    run.font.color.rgb = COLOR_TEXT
                    if ri <= 1:
                        run.font.bold = True
            # Couleur de fond pour les 2 premieres lignes (en-tete)
            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_TITLE
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                        run.font.size = Pt(14)
            elif ri == 1:
                cell.fill.solid()
                if ci == 1:
                    cell.fill.fore_color.rgb = RGBColor(0xE8, 0xF1, 0xFA)
                elif ci == 2:
                    cell.fill.fore_color.rgb = RGBColor(0xFD, 0xEE, 0xE8)
                else:
                    cell.fill.fore_color.rgb = COLOR_BG_LIGHT

    # Note de bas de slide
    add_text(s,
             "Note : le total Architecture 2 additionne le temps du "
             "pipeline mainline (qui a deja produit sections.json) et le "
             "temps du benchmark orchestrateur. C'est le mode de calcul "
             "demande (l'orchestrateur reutilise les sections deja "
             "generees pendant la reunion).",
             0.4, 6.75, 12.5, 0.6, size=11, color=COLOR_SUBTLE)


# ============================================================================
# Main
# ============================================================================

def main() -> None:
    prs = make_prs()

    slide_cover(prs)

    # Architecture 1
    slide_archi1_vue(prs)
    slide_archi1_agent(
        prs, n=1, label="RESUME (par chunk)", color=COLOR_ACCENT_1,
        role="Pour 1 chunk donne, produit titre + contexte + points "
             "factuels. Appele 10 fois (une par chunk).",
        prompt=PROMPT_RESUME,
        inp="texte du chunk = transcript brut\n"
            "(segments [HH:MM:SS] Speaker: texte\n"
            " concatenes pour ce chunk)",
        out='{\n  "titre": "Tour de table sur l\'IA...",\n'
            '  "contexte": "Les participants echangent...",\n'
            '  "points": [\n'
            '    "Mathieu presente son role...",\n'
            '    "Maya travaille sur des projets...",\n'
            '    ...\n'
            '  ]\n}',
    )
    slide_archi1_agent(
        prs, n=2, label="EXTRACTION decisions (par chunk)", color=COLOR_ACCENT_1,
        role="Pour 1 chunk, extrait uniquement les decisions formellement "
             "actees. Appele 10 fois.",
        prompt=PROMPT_EXTRACTION,
        inp="texte du chunk = transcript brut",
        out='{\n  "decisions": []\n}\n\n'
            "(ou liste si decisions explicites, ex.\n"
            ' ["Passage en production valide..."])',
    )
    slide_archi1_agent(
        prs, n=3, label="PLAN per chunk (par chunk)", color=COLOR_ACCENT_1,
        role="Pour 1 chunk, extrait 0 a 2 actions post-reunion typees "
             "(engagement ou suggestion). Appele 10 fois.",
        prompt=PROMPT_PLAN_PERCHUNK,
        inp="texte du chunk = transcript brut",
        out='{\n  "plan": [\n'
            '    {"action": "Iterer sur la boucle...",\n'
            '     "responsable": "Basile ou Gerald",\n'
            '     "echeance": "—",\n'
            '     "type": "engagement"},\n'
            '    {"action": "Evaluer la generalisation...",\n'
            '     "responsable": "—",\n'
            '     "echeance": "—",\n'
            '     "type": "suggestion"}\n'
            '  ]\n}',
    )
    slide_archi1_agent(
        prs, n=4, label="EXECUTIVE SUMMARY (1 appel final)",
        color=COLOR_ACCENT_1,
        role="Apres la reunion, lit l'intro du transcript + les "
             "titres/resumes de toutes les sections, et redige UN "
             "paragraphe de 5 phrases.",
        prompt=PROMPT_EXEC_SUMMARY,
        inp="intro = 50 premieres lignes du transcript\n"
            "contenu = '- {titre}: {resume}'\n"
            "          pour chaque section (1..10)",
        out="(texte libre, 1 paragraphe de 5 phrases)\n\n"
            '"La reunion vise a etablir une collaboration...\n'
            ' Les echanges portent sur des besoins techniques\n'
            ' en NLP et MCP, ainsi que sur l\'exploration..."',
    )
    slide_archi1_assemblage(prs)

    # Architecture 2
    slide_archi2_vue(prs)
    slide_archi2_orchestrateur(prs)
    slide_archi2_worker(prs)
    slide_archi2_renderer(prs)

    # Comparatif
    slide_comparatif(prs)

    out = Path(__file__).parent / "Architectures_CR_comparaison.pptx"
    prs.save(out)
    print(f"PPTX ecrit : {out}")


if __name__ == "__main__":
    main()
