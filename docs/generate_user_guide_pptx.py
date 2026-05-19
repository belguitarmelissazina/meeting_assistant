# -*- coding: utf-8 -*-
"""
Génère le support de formation (mode opératoire) « Meeting Assistant »
au format PowerPoint (.pptx), en français, pour des utilisateurs non techniques.

Le script installe python-pptx si nécessaire, puis écrit :
    docs/Guide_Utilisateur_Meeting_Assistant.pptx

Les captures d'écran ne sont PAS insérées : un encadré « 📷 CAPTURE À INSÉRER »
réserve l'emplacement et décrit précisément la capture à coller.

Usage :
    python docs/generate_user_guide_pptx.py
"""

import os
import subprocess
import sys

# --------------------------------------------------------------------------- #
#  Dépendance : python-pptx (auto-installation si absente)
# --------------------------------------------------------------------------- #
try:
    import pptx  # noqa: F401
except ImportError:
    print("python-pptx absent — installation…")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# --------------------------------------------------------------------------- #
#  Charte graphique
# --------------------------------------------------------------------------- #
BRAND        = RGBColor(0x1E, 0x40, 0xAF)   # bleu profond
BRAND_LIGHT  = RGBColor(0x3B, 0x82, 0xF6)   # bleu clair
DARK         = RGBColor(0x0F, 0x17, 0x2A)   # texte principal
MUTED        = RGBColor(0x55, 0x65, 0x77)   # texte secondaire
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
BG_SOFT      = RGBColor(0xF8, 0xFA, 0xFC)   # fond très clair
CAP_FILL     = RGBColor(0xEF, 0xF4, 0xFB)   # fond encadré capture
CAP_BORDER   = RGBColor(0x60, 0x7D, 0xA8)   # bordure encadré capture
ACCENT       = RGBColor(0xD9, 0x77, 0x06)   # ambre (conseils / attention)

FONT = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# --------------------------------------------------------------------------- #
#  Helpers bas niveau
# --------------------------------------------------------------------------- #
def _solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _no_fill(shape):
    shape.fill.background()
    shape.line.fill.background()


def _set_dashed_border(shape, color, width_pt=1.5):
    """Bordure pointillée (manipulation XML pour compatibilité maximale)."""
    ln = shape.line
    ln.color.rgb = color
    ln.width = Pt(width_pt)
    lnEl = shape.line._get_or_add_ln()
    for tag in ("a:prstDash",):
        for el in lnEl.findall(qn(tag)):
            lnEl.remove(el)
    dash = lnEl.makeelement(qn("a:prstDash"), {"val": "dash"})
    lnEl.append(dash)


def _txt(slide, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return box, tf


def _para(tf, text, size, color, bold=False, first=False,
          align=PP_ALIGN.LEFT, space_after=6, italic=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(0)
    r = p.add_run()
    r.text = text
    f = r.font
    f.name = FONT
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    return p


def _indent(p, inches):
    """Indentation gauche d'un paragraphe (via XML, marL en EMU)."""
    pPr = p._p.get_or_add_pPr()
    pPr.set("marL", str(int(Inches(inches))))
    pPr.set("indent", "0")


def _rect(slide, left, top, width, height, color, rounded=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        left, top, width, height)
    _solid(shp, color)
    shp.shadow.inherit = False
    return shp


def _bg(slide, color=WHITE):
    shp = _rect(slide, 0, 0, SLIDE_W, SLIDE_H, color)
    slide.shapes._spTree.remove(shp._element)
    slide.shapes._spTree.insert(2, shp._element)
    return shp


# --------------------------------------------------------------------------- #
#  Composition de slides
# --------------------------------------------------------------------------- #
def add_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_footer(slide, idx, total):
    _rect(slide, 0, SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06), BRAND)
    box, tf = _txt(slide, Inches(0.5), SLIDE_H - Inches(0.46),
                   Inches(8), Inches(0.3))
    _para(tf, "Meeting Assistant — Guide d'installation et d'utilisation",
          9, MUTED, first=True)
    box, tf = _txt(slide, SLIDE_W - Inches(1.6), SLIDE_H - Inches(0.46),
                   Inches(1.1), Inches(0.3))
    _para(tf, f"{idx} / {total}", 9, MUTED, first=True, align=PP_ALIGN.RIGHT)


def cover_slide(prs):
    s = add_blank(prs)
    _bg(s, BRAND)
    _rect(s, 0, SLIDE_H - Inches(2.2), SLIDE_W, Inches(2.2),
          RGBColor(0x16, 0x33, 0x8C))
    box, tf = _txt(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(0.5))
    _para(tf, "MODE OPÉRATOIRE", 18, BRAND_LIGHT, bold=True, first=True)
    box, tf = _txt(s, Inches(0.9), Inches(2.1), Inches(11.5), Inches(2.0))
    _para(tf, "Meeting Assistant", 54, WHITE, bold=True, first=True)
    _para(tf, "Installer et utiliser l'application — guide pas à pas",
          24, RGBColor(0xCB, 0xD9, 0xF5))
    box, tf = _txt(s, Inches(0.9), SLIDE_H - Inches(1.5), Inches(11.5),
                   Inches(1.0))
    _para(tf, "Enregistrez vos réunions · Transcription automatique · "
              "Compte rendu Word généré pour vous",
          16, RGBColor(0xCB, 0xD9, 0xF5), first=True)
    _para(tf, "Aucune compétence technique requise", 14,
          RGBColor(0x9F, 0xB8, 0xEC), italic=True)
    return s


def section_slide(prs, number, title, subtitle):
    s = add_blank(prs)
    _bg(s, BG_SOFT)
    _rect(s, 0, 0, Inches(0.35), SLIDE_H, BRAND)
    box, tf = _txt(s, Inches(1.1), Inches(2.6), Inches(11), Inches(0.6))
    _para(tf, f"SECTION {number}", 20, BRAND_LIGHT, bold=True, first=True)
    box, tf = _txt(s, Inches(1.1), Inches(3.2), Inches(11), Inches(1.4))
    _para(tf, title, 40, DARK, bold=True, first=True)
    if subtitle:
        box, tf = _txt(s, Inches(1.1), Inches(4.5), Inches(10.5), Inches(0.8))
        _para(tf, subtitle, 18, MUTED, first=True)
    return s


def _capture_box(slide, left, top, width, height, instruction):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = CAP_FILL
    box.shadow.inherit = False
    _set_dashed_border(box, CAP_BORDER, 1.75)

    _, tf = _txt(slide, left + Inches(0.3), top + Inches(0.3),
                 width - Inches(0.6), height - Inches(0.6),
                 anchor=MSO_ANCHOR.MIDDLE)
    _para(tf, "📷  CAPTURE À INSÉRER", 15, BRAND, bold=True, first=True,
          align=PP_ALIGN.CENTER, space_after=10)
    _para(tf, instruction, 13, DARK, align=PP_ALIGN.CENTER, space_after=10)
    _para(tf, "→ Collez ici votre capture d'écran (elle remplacera ce cadre).",
          11, MUTED, align=PP_ALIGN.CENTER, italic=True)


def content_slide(prs, idx, total, kicker, title, bullets,
                   capture=None, note=None):
    s = add_blank(prs)
    _bg(s, WHITE)

    # En-tête
    _rect(s, 0, 0, SLIDE_W, Inches(1.35), BRAND)
    box, tf = _txt(s, Inches(0.6), Inches(0.22), Inches(11.8), Inches(0.35))
    _para(tf, kicker.upper(), 13, RGBColor(0xBE, 0xCF, 0xF2), bold=True,
          first=True)
    box, tf = _txt(s, Inches(0.6), Inches(0.55), Inches(12.1), Inches(0.75))
    _para(tf, title, 27, WHITE, bold=True, first=True)

    body_top = Inches(1.7)
    body_h = SLIDE_H - Inches(2.5)

    if capture:
        text_w = Inches(6.2)
        cap_left = Inches(7.05)
        cap_w = SLIDE_W - cap_left - Inches(0.5)
    else:
        text_w = SLIDE_W - Inches(1.2)
        cap_left = None

    # Corps : puces
    _, tf = _txt(s, Inches(0.6), body_top, text_w, body_h)
    first = True
    for b in bullets:
        if isinstance(b, tuple):
            text, lvl = b
        else:
            text, lvl = b, 0
        if lvl == 0:
            _para(tf, "●  " + text, 16, DARK, bold=False, first=first,
                  space_after=10)
        else:
            p = _para(tf, "–  " + text, 14, MUTED, first=first,
                      space_after=7)
            _indent(p, 0.35)
        first = False

    if note:
        ny = SLIDE_H - Inches(0.95)
        bar = _rect(s, Inches(0.6), ny, text_w, Inches(0.5), RGBColor(
            0xFE, 0xF3, 0xC7), rounded=True)
        bar.line.color.rgb = ACCENT
        bar.line.width = Pt(0.75)
        _, tf = _txt(s, Inches(0.8), ny + Inches(0.06), text_w - Inches(0.4),
                     Inches(0.4), anchor=MSO_ANCHOR.MIDDLE)
        _para(tf, "💡  " + note, 12, RGBColor(0x92, 0x5A, 0x05), bold=True,
              first=True)

    if capture:
        _capture_box(s, cap_left, body_top, cap_w, body_h, capture)

    add_footer(s, idx, total)
    return s


def table_slide(prs, idx, total, kicker, title, rows, col_titles):
    s = add_blank(prs)
    _bg(s, WHITE)
    _rect(s, 0, 0, SLIDE_W, Inches(1.35), BRAND)
    box, tf = _txt(s, Inches(0.6), Inches(0.22), Inches(11.8), Inches(0.35))
    _para(tf, kicker.upper(), 13, RGBColor(0xBE, 0xCF, 0xF2), bold=True,
          first=True)
    box, tf = _txt(s, Inches(0.6), Inches(0.55), Inches(12.1), Inches(0.75))
    _para(tf, title, 27, WHITE, bold=True, first=True)

    n = len(rows) + 1
    tbl_h = Inches(5.0)
    gtbl = s.shapes.add_table(n, 2, Inches(0.6), Inches(1.7),
                              Inches(12.1), tbl_h).table
    gtbl.columns[0].width = Inches(3.8)
    gtbl.columns[1].width = Inches(8.3)

    for j, ct in enumerate(col_titles):
        c = gtbl.cell(0, j)
        c.text = ct
        c.fill.solid()
        c.fill.fore_color.rgb = BRAND
        para = c.text_frame.paragraphs[0]
        para.runs[0].font.bold = True
        para.runs[0].font.color.rgb = WHITE
        para.runs[0].font.size = Pt(15)
        para.runs[0].font.name = FONT

    for i, (a, b) in enumerate(rows, start=1):
        for j, val in enumerate((a, b)):
            c = gtbl.cell(i, j)
            c.text = val
            c.fill.solid()
            c.fill.fore_color.rgb = BG_SOFT if i % 2 else WHITE
            pr = c.text_frame.paragraphs[0]
            pr.runs[0].font.size = Pt(13)
            pr.runs[0].font.name = FONT
            pr.runs[0].font.color.rgb = DARK
            if j == 0:
                pr.runs[0].font.bold = True
            c.text_frame.word_wrap = True

    add_footer(s, idx, total)
    return s


def closing_slide(prs):
    s = add_blank(prs)
    _bg(s, BRAND)
    box, tf = _txt(s, Inches(0.9), Inches(2.4), Inches(11.5), Inches(1.0))
    _para(tf, "Vous êtes prêt(e) !", 44, WHITE, bold=True, first=True)
    box, tf = _txt(s, Inches(0.9), Inches(3.5), Inches(11.5), Inches(2.4))
    _para(tf, "Démarrage rapide :", 18, BRAND_LIGHT, bold=True, first=True)
    for t in [
        "1.  Installer l'application (installeur Windows).",
        "2.  Laisser le premier lancement télécharger les modèles (une seule fois).",
        "3.  Cliquer « Réunion hors agenda » ou choisir une réunion de l'agenda.",
        "4.  Enregistrer ou importer l'audio, renseigner le contexte.",
        "5.  « Lancer le traitement » et récupérer le compte rendu Word.",
    ]:
        _para(tf, t, 16, RGBColor(0xDD, 0xE7, 0xF8), space_after=8)
    return s


def main():
    here = os.path.dirname(os.path.abspath(__file__))
    out_path = os.path.join(here, "Guide_Utilisateur_Meeting_Assistant.pptx")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Couverture
    cover_slide(prs)

    # Définition ordonnée du contenu : (type, payload)
    plan = []

    plan.append(("section", (1, "Découvrir l'application",
                 "Comprendre ce que fait Meeting Assistant et "
                 "ce dont vous avez besoin")))
    plan.append(("content", dict(
        kicker="À quoi sert l'application",
        title="Un assistant qui rédige vos comptes rendus de réunion",
        bullets=[
            "Vous enregistrez la réunion (ou importez un audio / "
            "un transcript existant).",
            "L'application transcrit automatiquement la parole en texte.",
            "Elle identifie les différents intervenants "
            "(« qui a dit quoi » — la diarisation).",
            "Elle génère un compte rendu Word professionnel : résumé, "
            "décisions, actions à mener.",
            ("Tout se passe sur votre ordinateur — vos données ne "
             "partent pas sur Internet.", 1),
            ("Option : intégration de votre agenda Microsoft.", 1),
        ],
        capture=None,
        note="Gain de temps : plus besoin de prendre des notes ni de "
             "rédiger le compte rendu manuellement.")))
    plan.append(("content", dict(
        kicker="Avant de commencer",
        title="Ce dont vous avez besoin",
        bullets=[
            "Un PC sous Windows 10 ou 11 (64 bits).",
            "Environ 3 à 4 Go d'espace disque libre.",
            "Un microphone (intégré ou casque) pour l'enregistrement "
            "en direct.",
            "Une connexion Internet UNIQUEMENT au premier lancement.",
            ("Ensuite, l'application fonctionne hors ligne.", 1),
            "Aucun compte ni clé n'est nécessaire pour démarrer.",
        ],
        capture=None,
        note="Le premier lancement télécharge ~2,3 Go : prévoyez une "
             "bonne connexion et un peu de patience (une seule fois).")))

    plan.append(("section", (2, "Installer l'application",
                 "De l'installeur au premier démarrage — 6 étapes")))
    plan.append(("content", dict(
        kicker="Installation — étape 1/6",
        title="Télécharger l'installeur",
        bullets=[
            "Récupérez le fichier d'installation fourni :",
            ("« Meeting Assistant-Setup-0.1.0.exe »", 1),
            "Enregistrez-le dans un dossier facile à retrouver "
            "(ex. Téléchargements).",
            "Le fichier est volumineux (~3 Go) : le téléchargement "
            "peut prendre quelques minutes.",
        ],
        capture="Capture du lien / e-mail / page d'où l'on télécharge "
                "le fichier « Meeting Assistant-Setup-0.1.0.exe » "
                "(ou du fichier dans Téléchargements).")))
    plan.append(("content", dict(
        kicker="Installation — étape 2/6",
        title="Lancer l'installeur (avertissement Windows)",
        bullets=[
            "Double-cliquez sur le fichier téléchargé.",
            "Windows peut afficher « Windows a protégé votre "
            "ordinateur » (SmartScreen).",
            "Cliquez sur « Informations complémentaires ».",
            "Puis sur « Exécuter quand même ».",
            ("C'est normal pour une application récente : "
             "le fichier n'est pas dangereux.", 1),
        ],
        capture="Capture de la fenêtre bleue Windows SmartScreen, "
                "avec le bouton « Exécuter quand même » visible.")))
    plan.append(("content", dict(
        kicker="Installation — étape 3/6",
        title="Assistant d'installation",
        bullets=[
            "L'assistant d'installation s'ouvre.",
            "Vous pouvez conserver le dossier proposé par défaut.",
            "Aucun mot de passe administrateur n'est nécessaire.",
            "Cliquez sur « Installer » et patientez.",
        ],
        capture="Capture de l'écran de l'assistant d'installation "
                "(choix du dossier / bouton Installer).")))
    plan.append(("content", dict(
        kicker="Installation — étape 4/6",
        title="Fin de l'installation",
        bullets=[
            "Un message confirme que l'installation est terminée.",
            "Un raccourci « Meeting Assistant » est créé sur le "
            "Bureau et le menu Démarrer.",
            "Cliquez sur « Terminer ».",
        ],
        capture="Capture du dernier écran de l'assistant (« Terminer ») "
                "ET de l'icône Meeting Assistant sur le Bureau.")))
    plan.append(("content", dict(
        kicker="Installation — étape 5/6",
        title="Premier lancement : téléchargement des modèles IA",
        bullets=[
            "Double-cliquez sur l'icône « Meeting Assistant ».",
            "Une fenêtre « Préparation de Meeting Assistant » apparaît.",
            "L'application télécharge ~2,3 Go de modèles.",
            "Cela dure de 5 à 30 minutes selon votre connexion.",
            ("Ne fermez pas la fenêtre — cela n'arrive qu'une fois.", 1),
        ],
        capture="Capture de la fenêtre « Préparation de Meeting "
                "Assistant » avec la barre de progression.",
        note="Si le téléchargement est interrompu, il reprend tout "
             "seul là où il s'était arrêté.")))
    plan.append(("content", dict(
        kicker="Installation — étape 6/6",
        title="Démarrage de l'application",
        bullets=[
            "Une fois les modèles prêts, l'écran « Meeting Assistant — "
            "Démarrage… » s'affiche brièvement.",
            "L'application s'ouvre ensuite automatiquement.",
            "Les démarrages suivants seront beaucoup plus rapides.",
        ],
        capture="Capture de l'écran de démarrage "
                "« Meeting Assistant — Démarrage… ».")))

    plan.append(("section", (3, "Découvrir l'interface",
                 "Repérer les zones principales et les réglages "
                 "facultatifs")))
    plan.append(("content", dict(
        kicker="L'écran principal",
        title="Se repérer dans l'application",
        bullets=[
            "En haut à gauche : le menu ☰ (liste de vos réunions).",
            "Bouton « Réunion hors agenda » : réunion non planifiée.",
            "Icône Paramètres ⚙ : réglages facultatifs.",
            "Barre latérale : historique et dossiers.",
            "Zone centrale : agenda et comptes rendus.",
        ],
        capture="Capture de l'écran d'accueil. Entourez les 4 zones : "
                "menu ☰, bouton « Réunion hors agenda », "
                "Paramètres ⚙, barre latérale.")))
    plan.append(("content", dict(
        kicker="Réglage facultatif",
        title="Connecter votre agenda Microsoft (optionnel)",
        bullets=[
            "Retrouvez vos réunions Outlook dans l'application.",
            "Cliquez « Connecter mon agenda Microsoft ».",
            "Un code s'affiche : ouvrez microsoft.com/devicelogin "
            "et saisissez-le.",
            "Connectez-vous avec votre compte Microsoft habituel.",
            ("Le mot de passe n'est jamais stocké ; lecture seule "
             "de l'agenda.", 1),
        ],
        capture="Capture de l'écran « Connexion Microsoft » avec le "
                "code et le bouton « Ouvrir la page de connexion ».")))
    plan.append(("content", dict(
        kicker="Réglage facultatif",
        title="Compte rendu plus rapide : clé Mistral (optionnel)",
        bullets=[
            "Par défaut : compte rendu « Local » (hors ligne, "
            "gratuit, un peu plus lent).",
            "Option « Mistral Large » : plus rapide via Internet, "
            "nécessite une clé API.",
            "Paramètres ⚙ → « Clé API Mistral » → coller → "
            "« Enregistrer ».",
            ("Clé à obtenir sur console.mistral.ai (peut être payant).", 1),
            "Sans clé, gardez le mode « Local » : tout fonctionne.",
        ],
        capture="Capture de la boîte « Paramètres » montrant le champ "
                "« Clé API Mistral » et le bouton « Enregistrer ».")))

    plan.append(("section", (4, "Utiliser l'application",
                 "4 façons de capturer une réunion + le contexte")))
    plan.append(("content", dict(
        kicker="Cas d'usage 1",
        title="Enregistrer une réunion depuis l'agenda",
        bullets=[
            "Sur l'accueil, vos réunions à venir s'affichent.",
            "Cliquez sur la réunion concernée.",
            "Cliquez sur « Enregistrer ».",
            "Titre, participants et contexte sont pré-remplis "
            "automatiquement.",
        ],
        capture="Capture d'une carte de réunion à venir sur l'accueil "
                "avec le bouton « Enregistrer ».")))
    plan.append(("content", dict(
        kicker="Cas d'usage 2",
        title="Réunion hors agenda (enregistrement en direct)",
        bullets=[
            "Cliquez « Réunion hors agenda » (ou « Nouvelle réunion »).",
            "Écran « Capturez votre réunion » : onglet « Enregistrer ».",
            "Cliquez « Démarrer » : un minuteur défile.",
            "À la fin, cliquez « Arrêter ».",
        ],
        capture="Capture de l'écran « Capturez votre réunion » avec "
                "les onglets Enregistrer / Audio / Transcript et "
                "le minuteur.")))
    plan.append(("content", dict(
        kicker="Cas d'usage 3",
        title="Importer un fichier audio existant",
        bullets=[
            "« Réunion hors agenda » puis onglet « Audio ».",
            "Glissez votre fichier audio dans la zone "
            "(ou cliquez pour parcourir).",
            "Formats acceptés : audio courant et vidéo MP4.",
            "L'audio est traité comme un enregistrement.",
        ],
        capture="Capture de la zone « Glissez votre audio ici » "
                "(onglet Audio).")))
    plan.append(("content", dict(
        kicker="Cas d'usage 4",
        title="Importer un transcript déjà produit",
        bullets=[
            "« Réunion hors agenda » puis onglet « Transcript ».",
            "Glissez un .docx (transcript Teams) ou .txt formaté.",
            "Le texte est normalisé puis le compte rendu est généré.",
        ],
        capture="Capture de la zone « Glissez votre transcript ici » "
                "(onglet Transcript).")))
    plan.append(("content", dict(
        kicker="Préparer le traitement",
        title="Renseigner le contexte de la réunion",
        bullets=[
            "« Contexte » : sujet, enjeux, décisions attendues.",
            "« Participants » : noms séparés par des virgules.",
            "« Entreprises » : organisations impliquées.",
            "Réunion d'agenda : champs pré-remplis — vérifiez "
            "et ajustez.",
        ],
        capture="Capture du formulaire de contexte (Contexte, "
                "Participants, Entreprises).",
        note="Plus le contexte est précis, plus le compte rendu "
             "est pertinent.")))
    plan.append(("content", dict(
        kicker="Préparer le traitement",
        title="Choisir les options et lancer",
        bullets=[
            "Interrupteur « Diarisation » : identifie qui parle "
            "(recommandé).",
            "« Moteur de compte rendu » : « Local » ou "
            "« Mistral Large ».",
            "Mistral nécessite la clé API.",
            "Cliquez sur « Lancer le traitement ».",
        ],
        capture="Capture du bloc d'options : Diarisation, "
                "Local / Mistral Large, bouton "
                "« Lancer le traitement ».")))
    plan.append(("content", dict(
        kicker="Pendant le traitement",
        title="Suivre la progression",
        bullets=[
            "Une barre de progression affiche les étapes :",
            ("Conversion → Diarisation → Transcription → "
             "Compte rendu.", 1),
            "L'étape en cours est mise en évidence.",
            "Vous pouvez continuer à travailler en parallèle.",
            "Comptez ~1/5 à 1/3 de la durée de la réunion.",
        ],
        capture="Capture de la barre de progression du traitement "
                "(Conversion / Diarisation / Transcription / "
                "Compte rendu).")))

    plan.append(("section", (5, "Exploiter les résultats",
                 "Consulter, éditer, retrouver et organiser "
                 "vos comptes rendus")))
    plan.append(("content", dict(
        kicker="Le résultat",
        title="Consulter et modifier le compte rendu",
        bullets=[
            "Le compte rendu s'affiche dès qu'il est prêt.",
            "Vous pouvez le modifier directement dans l'application.",
            "Enregistrement automatique (« Enregistré à HH:MM »).",
            "Recherchez un mot avec Ctrl+F.",
        ],
        capture="Capture du compte rendu affiché (avec l'indicateur "
                "d'enregistrement automatique).")))
    plan.append(("content", dict(
        kicker="Le résultat",
        title="Écouter et télécharger",
        bullets=[
            "Bouton casque : ouvre un lecteur audio pour réécouter.",
            "Bouton de téléchargement : récupérer l'audio d'origine.",
            "Le compte rendu Word est aussi enregistré sur le disque.",
        ],
        capture="Capture du lecteur audio flottant (bas de l'écran) "
                "avec le bouton de téléchargement.")))
    plan.append(("content", dict(
        kicker="Le résultat",
        title="Où sont enregistrés mes fichiers ?",
        bullets=[
            "Documents ▸ Réunions ▸ (un dossier par réunion).",
            "« compte_rendu.docx » : le compte rendu Word.",
            "« transcript.txt » : la transcription complète.",
            "Le fichier audio de la réunion.",
            ("Gérez ces fichiers comme tout document "
             "(copier, envoyer…).", 1),
        ],
        capture="Capture de l'Explorateur Windows sur "
                "Documents\\Réunions montrant un dossier de "
                "réunion et ses fichiers.")))
    plan.append(("content", dict(
        kicker="S'organiser",
        title="Retrouver et classer vos réunions",
        bullets=[
            "Ouvrez la barre latérale avec le menu ☰.",
            "« Historique » : toutes vos réunions et leur statut.",
            "Bouton « Dossier » : créer des dossiers "
            "(client, projet…).",
            "Recherche (Ctrl+F) et mini-calendrier pour filtrer.",
        ],
        capture="Capture de la barre latérale ouverte : historique, "
                "dossiers, recherche, mini-calendrier.")))

    plan.append(("section", (6, "Maintenance & aide",
                 "Mises à jour, dépannage et bonnes pratiques")))
    plan.append(("content", dict(
        kicker="Mises à jour",
        title="L'application se met à jour toute seule",
        bullets=[
            "Message : « Une nouvelle version de Meeting Assistant "
            "est prête. »",
            "« Redémarrer maintenant » : installation immédiate.",
            "« Plus tard » : installation au prochain démarrage.",
            ("Aucune action technique de votre part.", 1),
        ],
        capture="Capture de la fenêtre de mise à jour "
                "« Une nouvelle version est prête » "
                "(boutons Redémarrer / Plus tard).")))
    plan.append(("table", dict(
        kicker="Aide",
        title="Questions fréquentes / dépannage",
        col_titles=["Situation", "Que faire"],
        rows=[
            ("Premier lancement très long",
             "Normal : téléchargement unique des modèles (~2,3 Go). "
             "Patientez, ne fermez pas la fenêtre."),
            ("Windows bloque l'installeur",
             "« Informations complémentaires » puis "
             "« Exécuter quand même » (SmartScreen)."),
            ("Pas d'Internet",
             "Après installation, l'application fonctionne hors "
             "ligne (mode Local)."),
            ("Le micro ne marche pas",
             "Autorisez le microphone dans Windows et vérifiez le "
             "bon périphérique."),
            ("Traitement long",
             "Comptez ~1/5 à 1/3 de la durée de réunion. Travail "
             "en parallèle possible."),
            ("Mistral non disponible",
             "Renseignez la clé API, ou utilisez le mode "
             "« Local »."),
        ])))
    plan.append(("content", dict(
        kicker="Aide",
        title="Bonnes pratiques pour un bon compte rendu",
        bullets=[
            "Bon microphone, peu de bruit de fond.",
            "Renseignez contexte et participants : résultat "
            "nettement meilleur.",
            "Relisez et ajustez le compte rendu (éditable).",
            "Activez la diarisation pour distinguer les "
            "intervenants.",
            "Besoin de rapidité ? Moteur « Mistral Large ».",
        ],
        capture=None,
        note="Le compte rendu est une aide à la rédaction : une "
             "relecture humaine reste recommandée.")))

    # total = cover + plan items + closing
    total = len(plan) + 2  # +cover +closing
    idx = 2  # cover = 1, premier élément de plan = 2
    for kind, payload in plan:
        if kind == "section":
            number, title, subtitle = payload
            section_slide(prs, number, title, subtitle)
        elif kind == "content":
            content_slide(prs, idx, total,
                           payload["kicker"], payload["title"],
                           payload["bullets"],
                           capture=payload.get("capture"),
                           note=payload.get("note"))
        elif kind == "table":
            table_slide(prs, idx, total,
                        payload["kicker"], payload["title"],
                        payload["rows"], payload["col_titles"])
        idx += 1

    closing_slide(prs)

    prs.save(out_path)
    n = len(prs.slides._sldIdLst)
    size_kb = os.path.getsize(out_path) // 1024
    print(f"OK — {n} slides écrites.")
    print(f"Fichier : {out_path} ({size_kb} Ko)")


if __name__ == "__main__":
    main()
