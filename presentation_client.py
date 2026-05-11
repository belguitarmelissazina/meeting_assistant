"""Présentation 2 slides — Meeting Assistant.

Slide 1 — Cartes uniformes pour Problème / Besoin / Solution
Slide 2 — Screenshot + cartes (compte rendu, local, parallèle, autres)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# ── Palette Natran ──────────────────────────────────────────────────────────
NAVY        = RGBColor(0x18, 0x39, 0x46)
DARK        = RGBColor(0x3E, 0x36, 0x31)
GRAY        = RGBColor(0x79, 0x93, 0x93)
GRAY_LINE   = RGBColor(0xD0, 0xD7, 0xDB)
GRAY_LIGHT  = RGBColor(0xE5, 0xE9, 0xEC)

CORAL       = RGBColor(0xB8, 0x52, 0x4E)
CORAL_PALE  = RGBColor(0xF5, 0xD2, 0xCD)

TEAL        = RGBColor(0x40, 0x96, 0x9A)
TEAL_PALE   = RGBColor(0xF0, 0xF7, 0xFA)

GREEN       = RGBColor(0x69, 0x97, 0x47)
GREEN_PALE  = RGBColor(0xED, 0xF4, 0xE8)

WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
PAGE_BG     = RGBColor(0xFF, 0xFF, 0xFF)

FONT_T = "Calibri Light"
FONT_B = "Calibri"


# ── Helpers ─────────────────────────────────────────────────────────────────
def _no_shadow(shape):
    shape.shadow.inherit = False


def add_rect(slide, l, t, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(0.5)
    _no_shadow(s)
    return s


def add_round(slide, l, t, w, h, fill, line=None, adj=0.04):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.adjustments[0] = adj
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(0.5)
    _no_shadow(s)
    return s


def add_oval(slide, l, t, w, h, fill, line=None, line_w=0.75):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    _no_shadow(s)
    return s


def add_line(slide, l, t, w, h, color, weight=0.75):
    s = slide.shapes.add_connector(1, l, t, l + w, t + h)
    s.line.color.rgb = color
    s.line.width = Pt(weight)
    return s


def add_text(slide, l, t, w, h, text, *,
             size=12, bold=False, italic=False, color=DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT_B):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return tb


def add_footer(slide, page, total):
    add_line(slide, Inches(0.55), Inches(7.10),
             Inches(12.23), 0, GRAY_LINE, weight=0.5)
    add_text(slide, Inches(0.55), Inches(7.18), Inches(8), Inches(0.22),
             "Meeting Assistant  ·  Présentation client",
             size=9, color=GRAY, font=FONT_T)
    add_text(slide, Inches(8), Inches(7.18), Inches(4.78), Inches(0.22),
             f"{page} / {total}",
             size=9, color=GRAY, font=FONT_T, align=PP_ALIGN.RIGHT)


def add_page_title(slide, kicker, title, kicker_color=TEAL):
    add_text(slide, Inches(0.55), Inches(0.42), Inches(12), Inches(0.30),
             kicker.upper(), size=10, bold=True, color=kicker_color, font=FONT_T)
    add_text(slide, Inches(0.55), Inches(0.72), Inches(12.25), Inches(0.55),
             title, size=22, bold=True, color=NAVY, font=FONT_T)
    add_line(slide, Inches(0.55), Inches(1.35),
             Inches(12.23), 0, GRAY_LINE, weight=0.5)


def add_uniform_card(slide, l, t, w, h, *,
                     bg, accent, glyph, title, desc):
    """Carte uniforme avec icône cercle, titre et 1 ligne de description."""
    add_round(slide, l, t, w, h, bg, adj=0.05)
    # cercle icône en haut centré
    icon_d = Inches(0.65)
    cx = l + (w - icon_d) / 2
    cy = t + Inches(0.25)
    add_oval(slide, cx, cy, icon_d, icon_d, accent)
    tb_icon = slide.shapes.add_textbox(cx, cy, icon_d, icon_d)
    tf_icon = tb_icon.text_frame
    tf_icon.margin_left = Emu(0); tf_icon.margin_right = Emu(0)
    tf_icon.margin_top = Emu(0); tf_icon.margin_bottom = Emu(0)
    tf_icon.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf_icon.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = glyph
    r.font.name = FONT_T; r.font.size = Pt(20); r.font.bold = True
    r.font.color.rgb = WHITE

    # titre — 1 ligne
    add_text(slide, l + Inches(0.20), cy + icon_d + Inches(0.15),
             w - Inches(0.40), Inches(0.35),
             title, size=12, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, font=FONT_T)
    # description — 1 phrase courte
    add_text(slide, l + Inches(0.20), cy + icon_d + Inches(0.50),
             w - Inches(0.40), h - icon_d - Inches(1.0),
             desc, size=10, color=DARK,
             align=PP_ALIGN.CENTER, font=FONT_B)


# ── Présentation ────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
blank = prs.slide_layouts[6]


# ============================================================================
# SLIDE 1 — Style "Contexte" : titre + sections à filet vertical + icônes ronds
# ============================================================================
s1 = prs.slides.add_slide(blank)
add_rect(s1, 0, 0, SW, SH, PAGE_BG)

# Grand titre "Contexte" en coral
add_text(s1, Inches(0.55), Inches(0.30), Inches(8), Inches(0.60),
         "Contexte", size=26, bold=True, color=CORAL, font=FONT_T)


def add_section_label_with_bar(slide, l_in, t_in, label, bar_color=TEAL):
    """Petit filet vertical fin + label de section."""
    add_rect(slide, Inches(l_in), Inches(t_in + 0.04),
             Inches(0.04), Inches(0.26), bar_color)
    add_text(slide, Inches(l_in + 0.18), Inches(t_in),
             Inches(6.0), Inches(0.32),
             label, size=13, bold=True, color=DARK, font=FONT_T)


def add_circle_icon_card(slide, cx_in, top_in, glyph, head_kw, head_rest, body,
                          icon_size_in=0.95):
    """Icône ronde + caption courte sous l'icône, mot-clé MAJ gras."""
    icon_d = Inches(icon_size_in)
    icon_l = Inches(cx_in) - icon_d / 2
    icon_t = Inches(top_in)
    add_oval(slide, icon_l, icon_t, icon_d, icon_d, TEAL_PALE)
    inner_d = Inches(icon_size_in * 0.62)
    inner_l = Inches(cx_in) - inner_d / 2
    inner_t = icon_t + (icon_d - inner_d) / 2
    add_oval(slide, inner_l, inner_t, inner_d, inner_d, WHITE,
             line=TEAL, line_w=1.0)
    tb_g = slide.shapes.add_textbox(inner_l, inner_t, inner_d, inner_d)
    tf_g = tb_g.text_frame
    tf_g.margin_left = Emu(0); tf_g.margin_right = Emu(0)
    tf_g.margin_top = Emu(0); tf_g.margin_bottom = Emu(0)
    tf_g.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_g = tf_g.paragraphs[0]; p_g.alignment = PP_ALIGN.CENTER
    r_g = p_g.add_run(); r_g.text = glyph
    r_g.font.name = FONT_T; r_g.font.size = Pt(18); r_g.font.bold = True
    r_g.font.color.rgb = TEAL

    cap_t = icon_t + icon_d + Inches(0.14)
    cap_l = Inches(cx_in) - Inches(1.40)
    cap_w = Inches(2.80)
    tb_c = slide.shapes.add_textbox(cap_l, cap_t, cap_w, Inches(0.95))
    tf_c = tb_c.text_frame
    tf_c.word_wrap = True
    tf_c.margin_left = Emu(0); tf_c.margin_right = Emu(0)
    tf_c.margin_top = Emu(0); tf_c.margin_bottom = Emu(0)
    p_c = tf_c.paragraphs[0]
    p_c.alignment = PP_ALIGN.CENTER
    p_c.line_spacing = 1.18
    r1 = p_c.add_run(); r1.text = head_kw + " "
    r1.font.name = FONT_T; r1.font.size = Pt(10); r1.font.bold = True
    r1.font.color.rgb = NAVY
    if head_rest:
        r2 = p_c.add_run(); r2.text = head_rest
        r2.font.name = FONT_T; r2.font.size = Pt(10)
        r2.font.color.rgb = NAVY
    p_c2 = tf_c.add_paragraph()
    p_c2.alignment = PP_ALIGN.CENTER
    p_c2.line_spacing = 1.18
    r3 = p_c2.add_run(); r3.text = body
    r3.font.name = FONT_B; r3.font.size = Pt(9)
    r3.font.color.rgb = DARK


# ── ZONE 1 : LE PROBLÈME — 3 icônes en rangée ──────────────────────────────
add_section_label_with_bar(s1, l_in=0.55, t_in=1.05,
                            label="Le problème", bar_color=TEAL)

prob_data = [
    ("⏱", "RÉDACTION", "CHRONOPHAGE.",
     "1 à 2 heures de travail après chaque réunion."),
    ("✍", "PRISE DE NOTES", "AU DÉTRIMENT DE L'ÉCOUTE.",
     "Participer et noter en même temps fait rater des échanges."),
    ("☁", "OUTILS CLOUD", "UNIQUEMENT.",
     "Les outils du marché posent des questions de confidentialité."),
]
prob_t = 1.55
section_l = 0.55
section_w = 12.23
n_prob = len(prob_data)
prob_centers = [section_l + section_w * (i + 0.5) / n_prob for i in range(n_prob)]
for cx, (glyph, kw, rest, body) in zip(prob_centers, prob_data):
    add_circle_icon_card(s1, cx, prob_t, glyph, kw, rest, body)


# ── ZONE 2 : LE BESOIN — 2 icônes en rangée ────────────────────────────────
add_section_label_with_bar(s1, l_in=0.55, t_in=3.55,
                            label="Le besoin", bar_color=TEAL)

need_data = [
    ("✓", "COMPTE RENDU", "STRUCTURÉ.",
     "Un livrable prêt à diffuser dès la fin de la réunion."),
    ("⌧", "CONFIDENTIALITÉ", "DES CONTENUS.",
     "Les contenus traités ne sortent pas de l'environnement de l'utilisateur."),
]
need_t = 4.05
n_need = len(need_data)
need_centers = [section_l + section_w * (i + 0.5) / n_need for i in range(n_need)]
for cx, (glyph, kw, rest, body) in zip(need_centers, need_data):
    add_circle_icon_card(s1, cx, need_t, glyph, kw, rest, body)


# ── BANDEAU SOLUTION — beaucoup plus grand, pleine largeur ─────────────────
sol_t = Inches(5.95)
sol_h = Inches(1.10)
add_round(s1, Inches(0.55), sol_t, Inches(12.23), sol_h,
          GREEN_PALE, adj=0.05)
# filet vert vertical à gauche
add_rect(s1, Inches(0.55), sol_t, Inches(0.10), sol_h, GREEN)

# icône produit à gauche
sol_icon_d = Inches(0.70)
sol_icon_l = Inches(0.95)
sol_icon_t = sol_t + (sol_h - sol_icon_d) / 2
add_oval(s1, sol_icon_l, sol_icon_t, sol_icon_d, sol_icon_d, GREEN)
tb_si = s1.shapes.add_textbox(sol_icon_l, sol_icon_t, sol_icon_d, sol_icon_d)
tf_si = tb_si.text_frame
tf_si.margin_left = Emu(0); tf_si.margin_right = Emu(0)
tf_si.margin_top = Emu(0); tf_si.margin_bottom = Emu(0)
tf_si.vertical_anchor = MSO_ANCHOR.MIDDLE
p_si = tf_si.paragraphs[0]; p_si.alignment = PP_ALIGN.CENTER
r_si = p_si.add_run(); r_si.text = "▤"
r_si.font.name = FONT_T; r_si.font.size = Pt(22); r_si.font.bold = True
r_si.font.color.rgb = WHITE

# Texte
text_l = sol_icon_l + sol_icon_d + Inches(0.30)
text_w = Inches(12.23) - (text_l - Inches(0.55)) - Inches(0.30)

# Kicker "SOLUTION PROPOSÉE"
add_text(s1, text_l, sol_t + Inches(0.13), text_w, Inches(0.25),
         "SOLUTION PROPOSÉE",
         size=10, bold=True, color=GREEN, font=FONT_T)
# Nom produit (grand)
add_text(s1, text_l, sol_t + Inches(0.36), text_w, Inches(0.35),
         "Meeting Assistant  ·  application de bureau Windows",
         size=15, bold=True, color=NAVY, font=FONT_T)
# Description (taille augmentée)
add_text(s1, text_l, sol_t + Inches(0.72), text_w, Inches(0.40),
         "Capte la réunion, transcrit, identifie les intervenants et "
         "génère un compte rendu Word — sur le poste, sans connexion internet.",
         size=11, color=DARK, font=FONT_B)


# Footer minimal
add_text(s1, Inches(0.55), Inches(7.20), Inches(8), Inches(0.20),
         "Meeting Assistant", size=9, italic=True, color=GRAY, font=FONT_T)
add_text(s1, Inches(8), Inches(7.20), Inches(4.78), Inches(0.20),
         "1 / 2", size=9, color=GRAY, font=FONT_T, align=PP_ALIGN.RIGHT)


# ============================================================================
# SLIDE 2 — Screenshot + cartes structurées par ordre logique
# ============================================================================
s2 = prs.slides.add_slide(blank)
add_rect(s2, 0, 0, SW, SH, PAGE_BG)

add_page_title(
    s2,
    kicker="02  ·  Notre proposition",
    title="Meeting Assistant — un outil agentique de comptes rendus",
    kicker_color=GREEN,
)
add_text(s2, Inches(0.55), Inches(1.45), Inches(12.25), Inches(0.30),
         "Pipeline IA multi-étapes embarqué dans une application bureau Windows.",
         size=11, italic=True, color=GRAY, font=FONT_B)

# ── Bandeau de mots-clés (chips, sans nom de modèle) ────────────────────────
chips_t = Inches(1.85)
chips_h = Inches(0.36)
chips_data = [
    "OUTIL AGENTIQUE",
    "PIPELINE MULTI-ÉTAPES",
    "100 % LOCAL",
    "TEMPS RÉEL",
    "DIARISATION",
    "HORS LIGNE",
]
chip_gap = Inches(0.08)


def _measure_chip_w(text, char_w=0.075, pad=0.30):
    return Inches(len(text) * char_w + pad * 2)


chip_widths = [_measure_chip_w(t) for t in chips_data]
total_chip_w = sum(chip_widths, Emu(0)) + chip_gap * (len(chips_data) - 1)
start_chip_l = (SW - total_chip_w) / 2
cur_l = start_chip_l
for txt, cw in zip(chips_data, chip_widths):
    add_round(s2, cur_l, chips_t, cw, chips_h, GREEN_PALE, adj=0.45)
    add_text(s2, cur_l, chips_t, cw, chips_h, txt,
             size=9, bold=True, color=GREEN,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_T)
    cur_l = cur_l + cw + chip_gap

# ── ZONE 1 (top) : screenshot mockup (gauche) + carte compte rendu (droite) ─
top_t = Inches(2.50)
top_h = Inches(2.55)

# ───── Placeholder pour la capture d'écran (gauche) ────────────────────────
mock_l = Inches(0.55)
mock_w = Inches(6.20)

# Cadre simple — rectangle vide, fond très pâle, bord gris fin
add_round(s2, mock_l, top_t, mock_w, top_h,
          RGBColor(0xFA, 0xFB, 0xFC),
          line=GRAY_LINE, adj=0.02)

# Icône "image" au centre
icon_d = Inches(0.85)
icon_l = mock_l + (mock_w - icon_d) / 2
icon_t = top_t + (top_h - icon_d) / 2 - Inches(0.20)
add_oval(s2, icon_l, icon_t, icon_d, icon_d, GREEN_PALE)
tb_icon = s2.shapes.add_textbox(icon_l, icon_t, icon_d, icon_d)
tf_icon = tb_icon.text_frame
tf_icon.margin_left = Emu(0); tf_icon.margin_right = Emu(0)
tf_icon.margin_top = Emu(0); tf_icon.margin_bottom = Emu(0)
tf_icon.vertical_anchor = MSO_ANCHOR.MIDDLE
p_icon = tf_icon.paragraphs[0]; p_icon.alignment = PP_ALIGN.CENTER
r_icon = p_icon.add_run(); r_icon.text = "▦"
r_icon.font.name = FONT_T; r_icon.font.size = Pt(28); r_icon.font.bold = True
r_icon.font.color.rgb = GREEN

# Indication discrète
add_text(s2, mock_l, icon_t + icon_d + Inches(0.18),
         mock_w, Inches(0.30),
         "Capture d'écran de l'application",
         size=12, bold=True, color=NAVY,
         align=PP_ALIGN.CENTER, font=FONT_T)
add_text(s2, mock_l, icon_t + icon_d + Inches(0.50),
         mock_w, Inches(0.25),
         "Insérez ici votre image — clic-droit › Modifier l'image",
         size=9, italic=True, color=GRAY,
         align=PP_ALIGN.CENTER, font=FONT_B)


# ───── Carte « Le compte rendu généré » (droite, top zone) ─────────────────
cr_l = mock_l + mock_w + Inches(0.30)
cr_w = SW - cr_l - Inches(0.55)
cr_t = top_t
cr_h = top_h

add_round(s2, cr_l, cr_t, cr_w, cr_h, WHITE, line=GREEN, adj=0.03)
add_rect(s2, cr_l, cr_t, Inches(0.10), cr_h, GREEN)

add_text(s2, cr_l + Inches(0.30), cr_t + Inches(0.20),
         cr_w - Inches(0.50), Inches(0.30),
         "Le compte rendu généré",
         size=14, bold=True, color=NAVY, font=FONT_T)
add_text(s2, cr_l + Inches(0.30), cr_t + Inches(0.55),
         cr_w - Inches(0.50), Inches(0.30),
         "Document Word structuré, prêt à diffuser, organisé en quatre parties :",
         size=10, italic=True, color=GRAY, font=FONT_B)

cr_sections = [
    ("Synthèse exécutive",
     "Vue d'ensemble de la réunion en un paragraphe."),
    ("Sujets abordés",
     "Découpage automatique par thème (segmentation sémantique), points clés par sujet."),
    ("Décisions prises",
     "Tableau récapitulatif des décisions formellement actées."),
    ("Plan d'attaque",
     "Actions à mener avec, lorsqu'ils sont mentionnés, le responsable et l'échéance."),
]
list_t = cr_t + Inches(1.00)
list_h = Inches(0.36)
for i, (head, tail) in enumerate(cr_sections):
    iy = list_t + list_h * i
    # bullet rond vert
    bx = cr_l + Inches(0.35)
    by = iy + Inches(0.06)
    add_oval(s2, bx, by, Inches(0.10), Inches(0.10), GREEN)
    add_text(s2, cr_l + Inches(0.55), iy,
             cr_w - Inches(0.70), Inches(0.20),
             head, size=10.5, bold=True, color=NAVY, font=FONT_T)
    add_text(s2, cr_l + Inches(0.55), iy + Inches(0.18),
             cr_w - Inches(0.70), Inches(0.20),
             tail, size=9, color=DARK, font=FONT_B)


# ── ZONE 2 (mid) : 100% LOCAL + TRAITEMENT PARALLÈLE ───────────────────────
mid_t = Inches(5.20)
mid_h = Inches(1.20)
mid_l_total = Inches(0.55)
mid_w_total = Inches(12.23)
mid_gap = Inches(0.20)
mid_card_w = (mid_w_total - mid_gap) / 2

# Carte LOCAL
loc_l = mid_l_total
add_round(s2, loc_l, mid_t, mid_card_w, mid_h, WHITE,
          line=GRAY_LINE, adj=0.04)
add_rect(s2, loc_l, mid_t, Inches(0.08), mid_h, GREEN)
# icône
add_oval(s2, loc_l + Inches(0.25), mid_t + Inches(0.25),
         Inches(0.45), Inches(0.45), GREEN)
tb_l = s2.shapes.add_textbox(loc_l + Inches(0.25), mid_t + Inches(0.25),
                              Inches(0.45), Inches(0.45))
tf_l = tb_l.text_frame
tf_l.margin_left = Emu(0); tf_l.margin_right = Emu(0)
tf_l.margin_top = Emu(0); tf_l.margin_bottom = Emu(0)
tf_l.vertical_anchor = MSO_ANCHOR.MIDDLE
p_l = tf_l.paragraphs[0]; p_l.alignment = PP_ALIGN.CENTER
r_l = p_l.add_run(); r_l.text = "◆"
r_l.font.name = FONT_T; r_l.font.size = Pt(15); r_l.font.bold = True
r_l.font.color.rgb = WHITE

add_text(s2, loc_l + Inches(0.85), mid_t + Inches(0.18),
         mid_card_w - Inches(1.0), Inches(0.30),
         "Traitement 100 % local · option cloud",
         size=12.5, bold=True, color=NAVY, font=FONT_T)
# texte description avec mots-clés en couleur
tb_loc = s2.shapes.add_textbox(loc_l + Inches(0.85), mid_t + Inches(0.50),
                                mid_card_w - Inches(1.0), Inches(0.65))
tf_loc = tb_loc.text_frame
tf_loc.word_wrap = True
tf_loc.margin_left = Emu(0); tf_loc.margin_right = Emu(0)
tf_loc.margin_top = Emu(0); tf_loc.margin_bottom = Emu(0)
p_loc = tf_loc.paragraphs[0]
p_loc.line_spacing = 1.20
parts = [
    ("Modèle de langage", True),
    (", ", False),
    ("transcription", True),
    (" et ", False),
    ("diarisation", True),
    (" tournent sur le poste via ", False),
    ("llama.cpp", True),
    (". ", False),
    ("Confidentialité", True),
    (" assurée. Option ", False),
    ("API cloud", True),
    (" disponible si qualité maximale recherchée.", False),
]
for txt, key in parts:
    r = p_loc.add_run()
    r.text = txt
    r.font.name = FONT_B; r.font.size = Pt(9.5)
    if key:
        r.font.bold = True; r.font.color.rgb = GREEN
    else:
        r.font.color.rgb = DARK

# Carte PARALLÈLE
par_l = mid_l_total + mid_card_w + mid_gap
add_round(s2, par_l, mid_t, mid_card_w, mid_h, WHITE,
          line=GRAY_LINE, adj=0.04)
add_rect(s2, par_l, mid_t, Inches(0.08), mid_h, GREEN)
add_oval(s2, par_l + Inches(0.25), mid_t + Inches(0.25),
         Inches(0.45), Inches(0.45), GREEN)
tb_p = s2.shapes.add_textbox(par_l + Inches(0.25), mid_t + Inches(0.25),
                              Inches(0.45), Inches(0.45))
tf_p = tb_p.text_frame
tf_p.margin_left = Emu(0); tf_p.margin_right = Emu(0)
tf_p.margin_top = Emu(0); tf_p.margin_bottom = Emu(0)
tf_p.vertical_anchor = MSO_ANCHOR.MIDDLE
p_p = tf_p.paragraphs[0]; p_p.alignment = PP_ALIGN.CENTER
r_p = p_p.add_run(); r_p.text = "⚡"
r_p.font.name = FONT_T; r_p.font.size = Pt(15); r_p.font.bold = True
r_p.font.color.rgb = WHITE

add_text(s2, par_l + Inches(0.85), mid_t + Inches(0.18),
         mid_card_w - Inches(1.0), Inches(0.30),
         "Traitement temps réel et parallèle",
         size=12.5, bold=True, color=NAVY, font=FONT_T)
tb_par = s2.shapes.add_textbox(par_l + Inches(0.85), mid_t + Inches(0.50),
                                mid_card_w - Inches(1.0), Inches(0.65))
tf_par = tb_par.text_frame
tf_par.word_wrap = True
tf_par.margin_left = Emu(0); tf_par.margin_right = Emu(0)
tf_par.margin_top = Emu(0); tf_par.margin_bottom = Emu(0)
p_par = tf_par.paragraphs[0]
p_par.line_spacing = 1.20
parts2 = [
    ("Workers parallèles", True),
    (" : transcription, diarisation et génération s'exécutent ", False),
    ("en simultané pendant l'enregistrement", True),
    (". À la fin de la réunion, l'essentiel du compte rendu est déjà ", False),
    ("prêt", True),
    (".", False),
]
for txt, key in parts2:
    r = p_par.add_run()
    r.text = txt
    r.font.name = FONT_B; r.font.size = Pt(9.5)
    if key:
        r.font.bold = True; r.font.color.rgb = GREEN
    else:
        r.font.color.rgb = DARK


# ── ZONE 3 (bottom) : autres fonctionnalités en bandeau compact ─────────────
sec_t = Inches(6.55)
sec_h = Inches(0.50)
sec_l = Inches(0.55)
sec_w = Inches(12.23)

add_text(s2, sec_l, sec_t - Inches(0.25), sec_w, Inches(0.20),
         "AUTRES FONCTIONNALITÉS",
         size=9, bold=True, color=GREEN, font=FONT_T)

secondary = [
    ("☰", "Découpage par sujet",
     "structuration thématique automatique"),
    ("↗", "3 façons d'alimenter",
     "enregistrement / audio / Teams .docx"),
    ("✎", "Éditeur intégré",
     "retouches dans l'app + ré-export Word"),
    ("▣", "Stockage centralisé",
     "Documents/Réunions, accessible Explorateur"),
]
n_sec = len(secondary)
sec_gap = Inches(0.15)
col_sec = (sec_w - sec_gap * (n_sec - 1)) / n_sec
for i, (gl, ttl, det) in enumerate(secondary):
    cl = sec_l + (col_sec + sec_gap) * i
    icon_box = Inches(0.40)
    add_round(s2, cl, sec_t, icon_box, icon_box,
              GREEN_PALE, adj=0.25)
    add_text(s2, cl, sec_t, icon_box, icon_box,
             gl, size=12, bold=True, color=GREEN,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_T)
    add_text(s2, cl + Inches(0.50), sec_t - Inches(0.02),
             col_sec - Inches(0.55), Inches(0.25),
             ttl, size=10, bold=True, color=NAVY, font=FONT_T)
    add_text(s2, cl + Inches(0.50), sec_t + Inches(0.21),
             col_sec - Inches(0.55), Inches(0.25),
             det, size=8.5, italic=True, color=GRAY, font=FONT_B)

add_footer(s2, page=2, total=2)


# ── Sauvegarde ──────────────────────────────────────────────────────────────
out = "Meeting_Assistant_Presentation_v11.pptx"
prs.save(out)
print(f"OK - {out} genere ({len(prs.slides)} slides).")
